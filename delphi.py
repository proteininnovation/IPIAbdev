#!/usr/bin/env python3
# ══════════════════════════════════════════════════════════════════════════════
# DELPHI — Deep End-to-end Learning Platform for antibody Developability
#          with High Interpretability
#
# Module      : delphi.py
# Description : Main command-line interface for DELPHI. Supports embedding
#               generation, k-fold cross-validation, full-data training,
#               and prediction for binary antibody developability labels
#               (PSR polyreactivity, SEC aggregation, HIC, AC-SINS,
#               viscosity, expression, and any custom 1/0 label).
#               Integrates five model architectures:
#                 transformer_lm      — Dual-branch Transformer on PLM
#                                       embeddings (IgBERT, ABlang2,
#                                       AntiBERTa2, AntiBERTy). Best AUC.
#                 transformer_onehot  — Dual-branch Transformer on one-hot
#                                       sequences. No PLM required.
#                 rf                  — Random Forest with SHAP.
#                 xgboost             — XGBoost with SHAP.
#                 cnn                 — CNN on PLM embeddings.
#               Benchmark: mean AUC 0.957 (PSR), 0.934 (SEC).
#               Applicable to any binary antibody property label.
# Developer   : Hoan Nguyen, PhD
# Company     : Institute for Protein Innovation (IPI)
# Date        : 2026-05
# Version     : 1.0.0
# ══════════════════════════════════════════════════════════════════════════════
#
# ══════════════════════════════════════════════════════════════════════════════
# MODELS & PLMs SUPPORTED
# ══════════════════════════════════════════════════════════════════════════════
#
#   --model  transformer_lm      Transformer + PLM embeddings  (best AUC)
#            transformer_onehot  Transformer + one-hot sequences (no PLM)
#            rf                  Random Forest (fast, SHAP, mutagenesis)
#            xgboost             XGBoost
#            cnn                 CNN
#
#   --lm     biophysical         charge, pI, hydrophobicity, R/K/W counts (26d)
#            kmer                1-mer + 2-mer AA frequencies (~440d)
#            onehot              VH+VL one-hot position encoding
#            onehot_vh           VH one-hot only
#            onehot_cdr3         HCDR3 one-hot only
#            ablang              ABlang2 480-dim   pip install ablang2
#            antiberty           AntiBERTy 512-dim  pip install antiberty
#            antiberta2          AntiBERTa2 1024-dim pip install transformers
#            antiberta2-cssp     AntiBERTa2-CSSP 1024-dim
#            igbert              IgBERT 1024-dim
#
# ══════════════════════════════════════════════════════════════════════════════
# TRAINING MODES (transformer_lm)
# ══════════════════════════════════════════════════════════════════════════════
#
#   MODE 1 — Frozen embeddings  (DEFAULT — recommended)
#   ─────────────────────────────────────────────────────────────────────────
#   sequences → PLM (frozen) → pre-computed .emb.csv → classifier trains
#   PLM weights: never updated  |  Trained: ~200k classifier params only
#   Best for: n < 10,000  |  CPU  |  same domain as pretraining
#   Your results: ablang ρ_OVA = −0.66 on GDPa3 (n=80)
#
#   MODE 2 — PLM layer unfreezing  (--train --finetune_plm)
#   ─────────────────────────────────────────────────────────────────────────
#   sequences → PLM (top layers update via backprop) → classifier trains
#   PLM weights: top N layers updated  |  Trained: ~20M params
#   Best for: n > 20,000  |  GPU  |  domain-shifted sequences
#   No .emb.csv needed — sequences processed in batches during training
#
#   MODE 3 — LoRA  (--train --finetune_plm --peft lora)  [RECOMMENDED]
#   ─────────────────────────────────────────────────────────────────────────
#   sequences → PLM (W frozen, only LoRA A×B trained) → classifier trains
#   PLM weights: W NEVER changes, only low-rank A×B matrices (~400k params)
#   Best for: n > 1,000  |  CPU feasible  |  low forgetting risk
#   No .emb.csv needed — sequences processed in batches during training
#
#   LEVEL 2 — Collaborator fine-tune  (--finetune --pretrained path.pt)
#   ─────────────────────────────────────────────────────────────────────────
#   Load YOUR pretrained .pt → fine-tune classifier on THEIR small dataset
#   PLM: embedded in their local DELPHI install (ablang2/transformers)
#   Best for: 50–2,000 new antibodies from a collaborator lab
#
# ══════════════════════════════════════════════════════════════════════════════
# CDR3 CONVENTION (2026 UPDATE)
# ══════════════════════════════════════════════════════════════════════════════
#
#   CDR3 does NOT include the leading framework cysteine (Kabat/IMGT pos 104).
#   Correct CDR3 starts with A-R (Ala-Arg germline anchor), NOT C-A-R.
#
#   Example:
#     HSEQ CDR3 region:  ...C | ARGGPGYAVFDY | WG...
#                             ^--- CDR3 starts here (not at C)
#
#   strip_cdr3_c_prefix() removes the leading 'C' if present.
#   This replaces the old fix_cdr3_c_prefix() which incorrectly ADDED 'C'.
#
#   Impact on models:
#     - CDR3 position 1 = A (Ala)   ← first true CDR3 residue
#     - CDR3 position 2 = R (Arg)   ← germline anchor
#     - CDR3 position 3+= variable  ← developability-relevant positions
#
#   All transformer_onehot and RF/XGBoost models must be retrained after
#   this change as CDR3 sequences are shorter by 1 residue.
#
# ══════════════════════════════════════════════════════════════════════════════
# EXAMPLE 1 — RF biophysical (fastest, interpretable, no PLM needed)
# ══════════════════════════════════════════════════════════════════════════════
#
#   python delphi.py --kfold 10 \
#       --target psr_filter --lm biophysical --model rf \
#       --db data/ipi_psr_trainset.xlsx \
#       --cost_fn 3.0
#
#   python delphi.py --train \
#       --target psr_filter --lm biophysical --model rf \
#       --db data/ipi_psr_trainset.xlsx
#
#   python delphi.py --predict data/new_cohort.xlsx \
#       --target psr_filter --lm biophysical --model rf \
#       --db data/ipi_psr_trainset.xlsx \
#       --mutagenesis 50
#
#   # Predict with a pretrained model by direct path (no --db needed):
#   python delphi.py --predict data/new_cohort.xlsx \
#       --target psr_filter --lm biophysical --model rf \
#       --model_path pretrained_202605/FINAL_psr_filter_biophysical_rf_ipi_psr_trainset.pkl \
#       --mutagenesis 50
#
# ══════════════════════════════════════════════════════════════════════════════
# EXAMPLE 2 — transformer_onehot (VH+VL+CDR3 one-hot, no PLM)
# ══════════════════════════════════════════════════════════════════════════════
#
#   python delphi.py --kfold 10 \
#       --target psr_filter --lm onehot --model transformer_onehot \
#       --db data/ipi_psr_trainset.xlsx \
#       --cluster 0.8
#
#   python delphi.py --train \
#       --target psr_filter --lm onehot --model transformer_onehot \
#       --db data/ipi_psr_trainset.xlsx
#
# ══════════════════════════════════════════════════════════════════════════════
# ── Platform fixes — must be FIRST, before any imports ───────────────────────
import os
os.environ.setdefault("TOKENIZERS_PARALLELISM", "false")
os.environ.setdefault("OBJC_DISABLE_INITIALIZE_FORK_SAFETY", "YES")
# ─────────────────────────────────────────────────────────────────────────────
"""
IPI Antibody Developability Prediction Platform
Final Production Version — 2026
Supports: SEC & PSR | XGBoost & RF & CNN & Transformer (One-Hot) & Transformer (LM)

Changes (2026-05):
  * strip_cdr3_c_prefix(): removes leading framework C from CDR3 (replaces fix_cdr3_c_prefix)
    CDR3 = AR... not CAR... The C belongs to VH framework (Kabat pos 104).
  * Auto-extract CDR3 from HSEQ when CDR3 column is missing (ANARCI -> regex fallback)
  * RF model: fixed kfold/predict/train calls to pass X_df + embeddings separately
"""

from config import MODEL_DIR, PREDICTION_DIR
import argparse
import os
import re
import warnings
import pandas as pd
import numpy as np
import torch
from pathlib import Path

from embedding_generator import generate_embedding

import sys
import datetime


def _list_registered_models(target: str = None, lm: str = None,
                            model: str = None,
                            registry_path: str = "config/model_registry.yaml") -> None:
    """
    Print models registered in config/model_registry.yaml, optionally filtered
    by target / lm / model. Called by `delphi.py --list-models`.
    """
    import os as _os
    import yaml as _yaml

    if not _os.path.exists(registry_path):
        print(f"[list-models] No registry found at {registry_path}")
        print(f"[list-models] Train a model with --train to create it, or "
              f"download pretrained models with utils/download_zenodo.py")
        return

    try:
        with open(registry_path) as f:
            reg = _yaml.safe_load(f) or {}
    except Exception as e:
        print(f"[list-models] Could not read {registry_path}: {e}")
        return

    models = reg.get("models", {})
    if not models:
        print(f"[list-models] Registry {registry_path} contains no models.")
        return

    rows = []
    for mid, e in models.items():
        if target and e.get("target") != target:
            continue
        if lm and e.get("lm") != lm:
            continue
        if model and e.get("model") != model:
            continue
        rows.append((mid, e))

    filt = []
    if target: filt.append(f"target={target}")
    if lm:     filt.append(f"lm={lm}")
    if model:  filt.append(f"model={model}")
    filt_s = ("  [filter: " + ", ".join(filt) + "]") if filt else ""

    print(f"\n  Registered models in {registry_path}{filt_s}")
    print(f"  {'-'*108}")
    print(f"  {'model_id':<60} {'target':<11} {'lm':<14} {'model':<18} {'type'}")
    print(f"  {'-'*108}")
    for mid, e in rows:
        print(f"  {mid:<60} {e.get('target',''):<11} {e.get('lm',''):<14} "
              f"{e.get('model',''):<18} {e.get('type','')}")
    print(f"  {'-'*108}")
    print(f"  {len(rows)} model(s)"
          + (f" of {len(models)} total" if len(rows) != len(models) else "")
          + "\n")


def _register_model(model_path: str, target: str, lm: str, model: str,
                    trainset: str, model_type_tag: str = "full_train",
                    registry_path: str = "config/model_registry.yaml",
                    **extra) -> None:
    """
    Add or update an entry in config/model_registry.yaml after --train.

    The registry key is the model filename (basename of model_path).
    Required fields: trainset, target, lm, model, model_path.
    Optional fields (trained_at, kfold_auc, threshold, epochs, notes, ...)
    are passed via **extra. Existing 'notes' is preserved on update.
    """
    import os as _os
    import yaml as _yaml

    # Integration/smoke tests train temporary models and must not append
    # machine-specific paths to the tracked project registry.
    if _os.environ.get("DELPHI_DISABLE_MODEL_REGISTRY", "").lower() in {
        "1", "true", "yes", "on"
    }:
        print("[registry] skipped (DELPHI_DISABLE_MODEL_REGISTRY=1)")
        return

    key = _os.path.basename(model_path)
    _os.makedirs(_os.path.dirname(registry_path) or ".", exist_ok=True)

    # Load existing registry (preserve all other entries)
    reg = {}
    if _os.path.exists(registry_path):
        try:
            with open(registry_path) as f:
                reg = _yaml.safe_load(f) or {}
        except Exception as e:
            print(f"[registry] WARNING: could not read {registry_path} ({e}); "
                  f"starting fresh")
            reg = {}
    if "models" not in reg or not isinstance(reg.get("models"), dict):
        reg["models"] = {}

    # Preserve a user-edited notes field if the entry already exists
    prior = reg["models"].get(key, {})
    entry = {
        "trainset":   str(trainset),
        "target":     target,
        "lm":         lm,
        "model":      model,
        "model_path": str(model_path),
        "type":       model_type_tag,
        "trained_at": datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
        "notes":      prior.get("notes", ""),
    }
    # Merge any extra metadata (kfold_auc, threshold, epochs, ...)
    for k, v in extra.items():
        if v is not None:
            entry[k] = v

    reg["models"][key] = entry

    try:
        with open(registry_path, "w") as f:
            _yaml.safe_dump(reg, f, default_flow_style=False,
                            sort_keys=False, allow_unicode=True)
        print(f"[registry] registered '{key}' → {registry_path}")
    except Exception as e:
        print(f"[registry] WARNING: failed to write {registry_path} ({e})")


def _parse_model_filename(model_path: str):
    """
    Parse a DELPHI-convention checkpoint filename into its components.

    Pattern: FINAL_{target}_{lm}_{model}_{db_stem}.{pt|pkl}

    Because target ('psr_filter'), lm ('transformer'... no, lm is e.g.
    'antiberta2-cssp') and model ('transformer_lm') can all contain
    underscores, we anchor on the known model and lm vocabularies rather than
    splitting blindly on '_'.

    Returns dict(target, lm, model, db_stem) or None if the name does not
    follow the convention.
    """
    import os as _os
    _VALID_MODELS = ["transformer_onehot", "transformer_lm", "xgboost", "rf", "cnn"]
    _VALID_LMS = ["antiberta2-cssp", "antiberta2", "antiberty", "ablang", "igbert",
                  "onehot_hcdr3", "onehot_cdr3", "onehot_vh", "onehot",
                  "biophysical", "kmer", "seq", "none"]

    name = _os.path.splitext(_os.path.basename(str(model_path)))[0]
    if not name.startswith("FINAL_"):
        return None
    body = name[len("FINAL_"):]   # {target}_{lm}_{model}_{db_stem}

    # 1) Find the model token (longest match first so transformer_* wins).
    model = None
    for m in _VALID_MODELS:
        tok = f"_{m}_"
        idx = body.find(tok)
        if idx != -1:
            model = m
            left = body[:idx]                       # {target}_{lm}
            db_stem = body[idx + len(tok):]         # {db_stem}
            break
    if model is None:
        return None

    # 2) From the left part, the lm is the known-LM token at its end.
    lm = None
    for cand in _VALID_LMS:
        if left == cand or left.endswith("_" + cand):
            lm = cand
            target = left[:len(left) - len(cand)].rstrip("_")
            break
    if lm is None or not target:
        return None

    return {"target": target, "lm": lm, "model": model, "db_stem": db_stem}


def _dbname_from_model_id(model_id_or_path: str, target: str = "",
                           lm: str = "", model: str = "") -> str:
    """
    Extract db_stem from model_id or model_path filename.

    Pattern: FINAL_{target}_{lm}_{model}_{db_stem}.pt/.pkl
    Examples:
      FINAL_psr_filter_ablang_transformer_lm_ipi_psr_trainset.pt  → ipi_psr_trainset
      FINAL_sec_filter_biophysical_rf_ipi_sec_5000.pkl             → ipi_sec_5000
      FINAL_psr_filter_onehot_transformer_onehot_DS1.pt            → DS1
    """
    import os as _os
    name = _os.path.splitext(_os.path.basename(str(model_id_or_path)))[0]
    # Try exact prefix match: FINAL_{target}_{lm}_{model}_
    prefix = f"FINAL_{target}_{lm}_{model}_"
    if name.startswith(prefix):
        return name[len(prefix):]
    # Fallback: strip FINAL_ and any known prefix parts
    if name.startswith("FINAL_"):
        name = name[6:]   # remove "FINAL_"
        # Remove target parts
        for part in (target or "").split("_"):
            if part and name.startswith(part + "_"):
                name = name[len(part)+1:]
        # Remove lm parts
        for part in (lm or "").replace("-","_").split("_"):
            if part and name.startswith(part + "_"):
                name = name[len(part)+1:]
        # Remove model parts
        for part in (model or "").split("_"):
            if part and name.startswith(part + "_"):
                name = name[len(part)+1:]
        if name:
            return name
    return "default"



class _Tee:
    def __init__(self, log_path: str, mode: str = 'w'):
        self._terminal = sys.__stdout__
        self._log = open(log_path, mode, buffering=1, encoding='utf-8')
        print(f"[log] Writing to: {log_path}" + (" (append)" if mode == 'a' else ""), flush=True)

    def write(self, msg):
        self._terminal.write(msg)
        self._log.write(msg)

    def flush(self):
        self._terminal.flush()
        self._log.flush()

    def close(self):
        if self._log and not self._log.closed:
            self._log.close()

    @property
    def encoding(self):
        return self._terminal.encoding

    def isatty(self):
        return False


def _setup_logging(args, db_path: str) -> str:
    ts      = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    _lm     = args.lm or "none"
    lm_tag  = _lm.replace(",", "_")
    _LM_NORM = {"onehot_cdr3": "onehot_hcdr3"}
    if args.lm is not None:
        args.lm  = _LM_NORM.get(args.lm, args.lm)
    lm_tag   = (args.lm or "none").replace(",", "_")
    db_stem = (Path(db_path).stem if db_path
             else _dbname_from_model_id(
                 getattr(args, "model_id", None) or
                 getattr(args, "model_path", None) or "default",
                 getattr(args, "target", ""),
                 getattr(args, "lm", ""),
                 getattr(args, "model", "")))

    if args.train:
        name     = f"train_{args.target}_{lm_tag}_{args.model}_{db_stem}_{ts}.log"
        log_path = os.path.join(MODEL_DIR, name)
        mode_w   = 'w'
    elif args.kfold:
        name     = f"kfold_{args.target}_{lm_tag}_{args.model}_{db_stem}_k{args.kfold}_{ts}.log"
        log_path = os.path.join(MODEL_DIR, name)
        mode_w   = 'w'
    elif args.predict:
        p        = Path(args.predict)
        name     = f"predict_{p.stem}_{args.target}_{lm_tag}_{args.model}_{db_stem}_{ts}.log"
        log_path = str(p.parent / name)
        mode_w   = 'w'
    elif args.build_embedding:
        p        = Path(args.build_embedding)
        name     = f"embedding_{p.stem}_{lm_tag}_{ts}.log"
        log_path = str(p.parent / name)
        mode_w   = 'w'
    elif getattr(args, 'split_dataset', False):
        name     = f"split_{Path(db_path).stem}_{ts}.log"
        log_path = str(Path(db_path).parent / name)
        mode_w   = 'w'
    else:
        log_path = f"ipi_{ts}.log"
        mode_w   = 'w'

    os.makedirs(os.path.dirname(log_path) or '.', exist_ok=True)
    sys.stdout = _Tee(log_path, mode=mode_w)
    ts_str = datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S')
    print(f"[log] {'='*58}")
    print(f"[log] Started  : {ts_str}")
    print(f"[log] Command  : {' '.join(sys.argv)}")
    print(f"[log] Platform : {sys.platform}  Python {sys.version.split()[0]}")
    print(f"[log] {'='*58}")
    print()
    return log_path


from models.xgboost import XGBoostModel
from models.randomforest import RandomForestModel
from models.cnn import CNNModel
from models.transformer_onehot import TransformerOneHotModel
from models.transformer_lm import TransformerLMModel

os.makedirs(MODEL_DIR, exist_ok=True)
os.makedirs(PREDICTION_DIR, exist_ok=True)


def get_default_db_path():
    data_dir = "data"
    if not os.path.exists(data_dir):
        return None
    files = [f for f in os.listdir(data_dir)
             if f.startswith("ipi_antibodydb") and f.endswith(".xlsx")]
    if not files:
        return None
    files.sort(key=lambda x: os.path.getmtime(os.path.join(data_dir, x)), reverse=True)
    return os.path.join(data_dir, files[0])


# ===========================================================================
# CDR3 EXTRACTION + PREFIX HANDLING
# ===========================================================================

def _extract_cdr3_single_anarci(hseq: str) -> str:
    try:
        with warnings.catch_warnings():
            warnings.simplefilter("ignore")
            from anarci import anarci as _anarci
            res = _anarci([("s", hseq)], scheme="imgt", output=False)
        numbered_seqs = res[0]
        if not numbered_seqs or numbered_seqs[0] is None:
            return ""
        domains = numbered_seqs[0]
        if not domains:
            return ""
        numbering, chain_type = domains[0]
        return "".join(
            aa for (pos, ins), aa in numbering
            if 105 <= pos <= 117 and aa != "-"
        )
    except Exception:
        return ""


def _extract_cdr3_regex(hseq: str) -> str:
    hseq = hseq.upper().replace("-", "").replace(" ", "")
    for pat in [r"C([A-Z]{3,30}?)WG[A-Z]G", r"C([A-Z]{3,35}?)WG"]:
        m = list(re.finditer(pat, hseq))
        if m:
            return m[-1].group(1)
    return ""


def extract_cdr3_from_hseq(hseq_series: pd.Series, verbose: bool = True) -> pd.Series:
    n = len(hseq_series)
    cdr3s = []
    n_anarci = n_regex = n_fail = 0

    if verbose:
        print(f"\n  [CDR3] CDR3 column not found — extracting from {n:,} HSEQ sequences ...")
        print(f"  [CDR3] Method: ANARCI (IMGT 105-117) with regex fallback")

    for hseq in hseq_series:
        if not isinstance(hseq, str) or len(hseq.strip()) < 20:
            cdr3s.append("")
            n_fail += 1
            continue
        cdr3 = _extract_cdr3_single_anarci(hseq)
        if cdr3:
            n_anarci += 1
        else:
            cdr3 = _extract_cdr3_regex(hseq)
            if cdr3:
                n_regex += 1
            else:
                n_fail += 1
        cdr3s.append(cdr3)

    if verbose:
        print(f"  [CDR3] Done: ANARCI={n_anarci:,}  regex={n_regex:,}  failed={n_fail:,}")
        if n_fail > 0:
            pct = 100 * n_fail / max(n, 1)
            print(f"  [CDR3] WARNING: {n_fail:,} sequences ({pct:.1f}%) could not be extracted — CDR3 set to ''.")
        if n_regex > 0:
            print(f"  [CDR3] NOTE: {n_regex:,} sequences used regex fallback.")

    return pd.Series(cdr3s, index=hseq_series.index, name="CDR3")


def _ensure_cdr3(df: pd.DataFrame, verbose: bool = True) -> pd.DataFrame:
    if "HSEQ" not in df.columns:
        if verbose:
            print("  [CDR3] WARNING: No HSEQ column — cannot extract CDR3.")
        return df

    df = df.copy()
    _nan_like = {"nan", "none", "null", "n/a", "na", ""}

    if "CDR3" not in df.columns:
        if verbose:
            print("  [CDR3] 'CDR3' column not found.")
        df["CDR3"] = ""
    else:
        df["CDR3"] = df["CDR3"].fillna("").astype(str)
        df.loc[df["CDR3"].str.strip().str.lower().isin(_nan_like), "CDR3"] = ""

    empty_mask = df["CDR3"].str.len() == 0
    n_missing  = empty_mask.sum()

    if n_missing > 0:
        if verbose and n_missing < len(df):
            print(f"  [CDR3] {n_missing:,} rows have empty CDR3 — extracting from HSEQ.")
        extracted = extract_cdr3_from_hseq(df.loc[empty_mask, "HSEQ"], verbose=verbose)
        df.loc[empty_mask, "CDR3"] = extracted.values

    return df


def strip_cdr3_c_prefix(df: pd.DataFrame, cdr3_col: str = "CDR3",
                        verbose: bool = True) -> pd.DataFrame:
    """
    Remove the leading 'C' (VH framework cysteine, Kabat/IMGT pos 104) from
    CDR3 sequences if present.

    The conserved cysteine is NOT part of CDR3. Correct CDR3 starts with A-R
    (Ala-Arg germline anchor), not C-A-R.

    Example:
        CARGGPGYAVFDY  ->  ARGGPGYAVFDY   (leading C removed)
        ARGGPGYAVFDY   ->  ARGGPGYAVFDY   (unchanged — no leading C)

    This function replaces the old fix_cdr3_c_prefix() which incorrectly
    ADDED a 'C' prefix. All models trained after this change will use CDR3
    sequences that are 1 residue shorter and correctly numbered:
        Position 1 = A (Ala)   — first true CDR3 residue
        Position 2 = R (Arg)   — conserved germline anchor
        Position 3+ = variable — developability-relevant positions
    """
    if cdr3_col not in df.columns:
        return df

    df   = df.copy()
    mask = (df[cdr3_col].str.len() > 1) & (df[cdr3_col].str.startswith("C"))
    n    = mask.sum()

    if n > 0 and verbose:
        print(f"\n  [CDR3-C] Removing framework 'C' from {n:,} CDR3 sequence(s).")
        print(f"  [CDR3-C] The leading C is the VH framework cysteine (Kabat pos 104), not CDR3.")
        print(f"  [CDR3-C] Correct CDR3 starts with AR, e.g. ARGGPGYAVFDY (not CARGGPGYAVFDY).")
        shown = 0
        for bc, row in df[mask].iterrows():
            old = row[cdr3_col]
            print(f"    BARCODE={bc}  '{old}'  ->  '{old[1:]}'")
            shown += 1
            if shown >= 5:
                break
        if n > 5:
            print(f"    ... and {n - 5:,} more")
    elif verbose and n == 0:
        print(f"  [CDR3-C] No leading 'C' found — CDR3 sequences already correctly formatted.")

    df.loc[mask, cdr3_col] = df.loc[mask, cdr3_col].str[1:]
    return df


def fix_cdr3_c_prefix(df: pd.DataFrame, cdr3_col: str = "CDR3",
                      verbose: bool = True) -> pd.DataFrame:
    """
    Deprecated — now calls strip_cdr3_c_prefix() instead.

    Old behaviour was to ADD a 'C' prefix if missing, which was incorrect.
    The leading C belongs to VH framework, not CDR3. This wrapper is kept
    for backward compatibility only — it now removes the C instead.
    """
    return strip_cdr3_c_prefix(df, cdr3_col=cdr3_col, verbose=verbose)


# ===========================================================================
# FILE I/O HELPERS
# ===========================================================================

def save_dataframe(df: pd.DataFrame, path: str) -> None:
    path = str(path)
    os.makedirs(os.path.dirname(path) or '.', exist_ok=True)
    if path.lower().endswith('.csv'):
        df.to_csv(path, index=False)
    else:
        df.to_excel(path, index=False)
    print(f"  Saved: {path}  ({len(df):,} rows)")


def read_dataframe(path: str) -> pd.DataFrame:
    if path.lower().endswith('.csv'):
        return pd.read_csv(path)
    return pd.read_excel(path)


def _split_embedding(emb_path, train_barcodes, val_barcodes, train_out, val_out):
    emb = pd.read_csv(emb_path, index_col=0)
    emb.index = emb.index.astype(str).str.strip()
    db_base = str(Path(emb_path).name)
    _parts  = db_base.split(".")
    emb_idx = next((i for i, p in enumerate(_parts) if p == "emb"), None)
    if emb_idx is None:
        print(f"  [emb-split] Cannot parse LM tag from {db_base} — skipping")
        return
    lm_tag = ".".join(_parts[emb_idx-1:])
    tr_out = f"{train_out}.{lm_tag}"
    va_out = f"{val_out}.{lm_tag}"
    tr_bcs = [b for b in train_barcodes.astype(str) if b in emb.index]
    va_bcs = [b for b in val_barcodes.astype(str)   if b in emb.index]
    n_missing_tr = len(train_barcodes) - len(tr_bcs)
    n_missing_va = len(val_barcodes)   - len(va_bcs)
    if n_missing_tr or n_missing_va:
        print(f"  [emb-split] WARNING: {n_missing_tr} train + {n_missing_va} val BARCODEs not found — skipped")
    emb.loc[tr_bcs].to_csv(tr_out)
    emb.loc[va_bcs].to_csv(va_out)
    print(f"  [emb-split] {os.path.basename(tr_out)}  ({len(tr_bcs):,} rows)")
    print(f"  [emb-split] {os.path.basename(va_out)}  ({len(va_bcs):,} rows)")


def _find_embedding_files(db_path: str) -> list:
    import glob
    return sorted(glob.glob(f"{db_path}.*.emb.csv"))


def _align_embedding(df: pd.DataFrame, embedding: pd.DataFrame, context: str = "") -> tuple:
    tag = f"[{context}] " if context else ""
    df = df.copy()
    df.index = df.index.astype(str).str.strip()
    embedding = embedding.copy()
    embedding.index      = embedding.index.astype(str).str.strip()
    embedding.index.name = "BARCODE"
    _sample = embedding.index[:5].tolist()
    if all(str(v).isdigit() for v in _sample):
        print(f"  {tag}WARNING: embedding index looks numeric ({_sample}) — expected BARCODE strings.")
    merged = df.join(embedding, how="inner")
    n_missing = len(df) - len(merged)
    if len(merged) == 0:
        raise ValueError(
            f"{tag}No overlapping BARCODEs between data ({len(df):,} rows) "
            f"and embedding ({len(embedding):,} rows)."
        )
    if n_missing > 0:
        print(f"  {tag}WARNING: {n_missing:,} / {len(df):,} rows ({n_missing/len(df):.1%}) have no embedding — excluded.")
    emb_cols = embedding.columns.tolist()
    return merged[emb_cols], merged.drop(columns=emb_cols)


def split_and_save(db_path: str, split: float = 0.8,
                   cluster_thresh: float = 0.8,
                   cluster_col: str = "CDR3",
                   label_col: str = "psr_filter") -> tuple:
    from sklearn.model_selection import StratifiedGroupKFold
    from utils.clustering import greedy_clustering_by_levenshtein

    df = read_dataframe(db_path)
    p  = Path(db_path)
    print(f"\n[split] {p.name}  ({len(df):,} rows)  split={split:.0%}/{1-split:.0%}  cluster_col={cluster_col}")

    _cc = cluster_col.upper()
    if _cc == 'HSEQ':  _cc = 'VH'
    if _cc == 'HVHVL': _cc = 'VHVL'

    col_map = {'CDR3': f'HCDR3_CLUSTER_{cluster_thresh}',
               'VH':   f'VH_CLUSTER_{cluster_thresh}',
               'VHVL': f'VHVL_CLUSTER_{cluster_thresh}'}
    seq_map = {'CDR3': 'CDR3', 'VH': 'HSEQ', 'VHVL': 'HSEQ'}
    grp_col = col_map.get(_cc, f'HCDR3_CLUSTER_{cluster_thresh}')
    seq_col = seq_map.get(_cc, 'CDR3')

    if grp_col not in df.columns:
        if seq_col not in df.columns:
            raise ValueError(f"Sequence column '{seq_col}' not found.")
        print(f"[split] Computing {grp_col} (threshold={cluster_thresh}) ...")
        seqs = (df['HSEQ'].fillna('').astype(str) + '_' +
                df['LSEQ'].fillna('').astype(str)).tolist() if _cc == 'VHVL' \
               else df[seq_col].fillna('').astype(str).tolist()
        df[grp_col] = greedy_clustering_by_levenshtein(seqs, cluster_thresh)
        n_clust = df[grp_col].nunique()
        print(f"[split] {n_clust:,} clusters  (mean {len(df)/n_clust:.1f} sequences/cluster)")
        try:
            _df_save = read_dataframe(db_path)
            _df_save[grp_col] = df[grp_col].values
            save_dataframe(_df_save, db_path)
            print(f"[split] Saved {grp_col} → {os.path.basename(db_path)}")
        except Exception as _e:
            print(f"[split] WARNING: could not save {grp_col} back to {os.path.basename(db_path)}: {_e}")
    else:
        print(f"[split] Using existing {grp_col}  ({df[grp_col].nunique():,} clusters)")

    _lc     = label_col if label_col in df.columns else None
    has_lbl = _lc is not None

    if has_lbl:
        _lbl_mask = df[_lc].notna()
        n_unl     = (~_lbl_mask).sum()
        if n_unl > 0:
            print(f"[split] {_lbl_mask.sum():,} labelled + {n_unl:,} unlabelled rows — ALL preserved")
    else:
        print(f"[split] WARNING: '{label_col}' not found — splitting by cluster only")
        _lbl_mask = pd.Series([True] * len(df), index=df.index)

    df_lbl  = df[_lbl_mask].copy()
    y_arr   = df_lbl[_lc].values.astype(int) if has_lbl else np.zeros(len(df_lbl), int)
    groups  = df_lbl[grp_col].values

    n_splits = max(2, round(1.0 / (1.0 - split)))
    sgkf = StratifiedGroupKFold(n_splits=n_splits, shuffle=True, random_state=42)
    best_split, best_diff = None, float('inf')
    for tr, va in sgkf.split(np.arange(len(y_arr)), y_arr, groups):
        diff = abs(y_arr[va].mean() - y_arr.mean()) if has_lbl else 0.0
        if diff < best_diff: best_diff, best_split = diff, (tr, va)

    tr_idx, va_idx = best_split
    tr_clusters = set(groups[tr_idx])
    va_clusters = set(groups[va_idx])
    leaked = tr_clusters & va_clusters
    print(f"[split] {'[WARN] ' + str(len(leaked)) + ' cluster(s) leaked' if leaked else '✓ No cluster leakage'}")

    if not _lbl_mask.all():
        df_unl = df[~_lbl_mask].copy()
        unl_in_train = df_unl[grp_col].apply(lambda g: g in tr_clusters or g not in va_clusters)
        df_unl_train = df_unl[unl_in_train]
        df_unl_val   = df_unl[~unl_in_train]
    else:
        df_unl_train = df_unl_val = pd.DataFrame(columns=df.columns)

    train_df = pd.concat([df_lbl.iloc[tr_idx], df_unl_train], ignore_index=True)
    val_df   = pd.concat([df_lbl.iloc[va_idx],   df_unl_val], ignore_index=True)

    _pos_tr = f"  pos={y_arr[tr_idx].mean():.1%}" if has_lbl else ""
    _pos_va = f"  pos={y_arr[va_idx].mean():.1%}" if has_lbl else ""
    print(f"[split] Train={len(train_df):,}{_pos_tr}  Val={len(val_df):,}{_pos_va}")

    ext       = p.suffix
    train_out = str(p.parent / f"{p.stem}_train{ext}")
    val_out   = str(p.parent / f"{p.stem}_val{ext}")
    save_dataframe(train_df.reset_index(drop=True), train_out)
    save_dataframe(val_df.reset_index(drop=True),   val_out)

    bc_col = 'BARCODE'
    if bc_col in df.columns:
        train_bcs = train_df[bc_col].astype(str)
        val_bcs   = val_df[bc_col].astype(str)
    else:
        train_bcs = pd.Index(train_df.index.astype(str))
        val_bcs   = pd.Index(val_df.index.astype(str))

    emb_files = _find_embedding_files(db_path)
    if emb_files:
        print(f"\n[split] Found {len(emb_files)} embedding file(s) — splitting by BARCODE ...")
        for emb_path in emb_files:
            lm_name = os.path.basename(emb_path).split('.emb.csv')[0].split('.')[-1]
            print(f"  [emb-split] {lm_name}: {os.path.basename(emb_path)}")
            try:
                _split_embedding(emb_path, train_bcs, val_bcs, train_out, val_out)
            except Exception as e:
                print(f"  [emb-split] WARNING: failed for {lm_name} — {e}")
    else:
        print(f"\n[split] No embedding files found for {p.name}")

    return train_out, val_out, train_df, val_df


# ===========================================================================
# LOAD DATA
# ===========================================================================

def load_data(db_path, lm="antiberta2", label_col="psr_filter"):
    print(f"\nLoading database : {os.path.basename(db_path)}")
    print(f"Target           : {label_col}")
    print(f"Embedding        : {lm}")

    df = pd.read_excel(db_path)
    print(f"Total rows       : {len(df):,}  |  columns: {len(df.columns)}")

    required = ["BARCODE", "HSEQ", "LSEQ", label_col]
    if lm in ("onehot", "onehot_vh"):
        required += ["CDR3"]

    if "HSEQ" in df.columns:
        df = _ensure_cdr3(df, verbose=True)
        df = strip_cdr3_c_prefix(df, verbose=True)   # ← strip C, not add C

    missing = [c for c in required if c not in df.columns]
    if missing:
        raise ValueError(f"Missing required columns: {missing}")

    df = df.dropna(subset=required).set_index("BARCODE")
    print(f"After dropna     : {len(df):,} rows")

    _SEQ_ONLY_MODES = {"onehot", "onehot_vh", "onehot_cdr3", "onehot_hcdr3",
                       "biophysical", "kmer", "seq", "none"}

    if lm in _SEQ_ONLY_MODES:
        X    = df[["HSEQ", "LSEQ", "CDR3"]]
        y    = df[label_col].values
        data = df.copy()
        if lm not in ("onehot", "onehot_vh"):
            print(f"[load_data] lm='{lm}' — sequence-only mode, no PLM embedding loaded")
    else:
        possible = [f"{db_path}.{lm}.emb.csv"]
        emb_file = next((f for f in possible if os.path.exists(f)), None)
        if not emb_file:
            print(f"Embedding not found -> generating {lm}...")
            emb_file = generate_embedding(db_path, lm=lm)
        if not emb_file or not os.path.exists(str(emb_file or '')):
            _SEQ_ONLY = {"biophysical","kmer","onehot","onehot_vh",
                         "onehot_cdr3","onehot_hcdr3","none","seq"}
            _PLM_LIST = "ablang | antiberty | antiberta2 | antiberta2-cssp | igbert | igt5 | abmap"
            raise ValueError(
                f"\n[load_data] Cannot load embedding for lm='{lm}'.\n"
                f"  Embedding file not found and generation failed.\n"
                f"  If you intended a sequence-only model, use one of:\n"
                f"    --lm biophysical  --lm kmer  --lm onehot\n"
                f"  If you intended a PLM, supported values are:\n"
                f"    {_PLM_LIST}\n"
                f"  Note: --model rf/xgboost/cnn is NOT a valid --lm value.\n"
                f"  Example: --lm igbert --model rf")
        print(f"Embedding file   : {emb_file}")
        print(f"Embedding size   : {os.path.getsize(emb_file) / 1024 / 1024:.1f} MB")
        embedding = pd.read_csv(emb_file, index_col=0)
        print(f"Embedding shape  : {embedding.shape[0]:,} samples × {embedding.shape[1]:,} dims")
        X, data = _align_embedding(df, embedding, context="load_data")
        y = data[label_col].values

    print(f"Samples loaded   : {len(y):,}")
    print(f"Target stats     : mean={y.mean():.4f}  std={y.std():.4f}  "
          f"min={y.min():.4f}  max={y.max():.4f}")
    return X, data, y


# ── Universal CDR3 mutagenesis ────────────────────────────────────────────────

def _run_cdr3_mutagenesis(
        data, model, model_type, lm, db_stem,
        input_file, target, test_target,
        embeddings_fn=None, n_override=None):
    import traceback as _tb
    import matplotlib.pyplot as plt
    import matplotlib.colors as _mc
    import numpy as np

    AMINO_ACIDS = 'ACDEFGHIKLMNPQRSTVWY'

    _cfg = getattr(model, 'config', {}) if hasattr(model, 'config') else {}
    _mut_cfg  = _cfg.get('mutagenesis', {})
    _fmt      = _mut_cfg.get('format',   'tiff')
    _dpi      = _mut_cfg.get('pub_dpi',   300)
    _make_ppt = _mut_cfg.get('make_ppt',  True)
    _max_n    = _mut_cfg.get('max_samples', 50)
    if n_override is not None:
        _max_n = len(data) if n_override == 0 else int(n_override)

    _path     = Path(input_file)
    _stem     = f"{_path.stem}_{target}_{lm}_{model_type}_{db_stem}"
    _mut_dir  = str(_path.parent / f"{_stem}_mutagenesis")
    import os; os.makedirs(_mut_dir, exist_ok=True)

    _label_col = test_target if (test_target and test_target in data.columns) \
                 else (target if target in data.columns else None)

    _n = min(len(data), _max_n)
    if _n < len(data):
        print(f"[Mutagenesis] Limited to first {_n}/{len(data)} antibodies")

    _needs_emb = (hasattr(model, 'fb_') and model.fb_ is not None and
                  model.fb_.feat_cfg.get('embedding', False))

    if _needs_emb:
        print(f"\n[Mutagenesis] WARNING: PLM mode '{lm}' — heatmap may be flat.")
        print(f"  Recommended: --lm biophysical or --lm kmer for meaningful mutagenesis.")
        try:
            from embedding_generator import EmbeddingGenerator as _EmbGen
            _emb_gen = _EmbGen(lm=lm)
            def _embeddings_fn(df_row):
                _hseq = str(df_row.iloc[0].get('HSEQ', '') or '')
                _lseq = str(df_row.iloc[0].get('LSEQ', '') or '')
                _bc   = str(df_row.index[0])
                _emb  = _emb_gen.embed_single(_hseq, _lseq, barcode=_bc)
                return _emb.reshape(1, -1).astype('float32')
        except Exception as _ee:
            try:
                import tempfile, os as _os, pandas as _pd
                from embedding_generator import generate_embedding as _gen_emb
                def _embeddings_fn(df_row):
                    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as _tf:
                        _tmp = _tf.name
                    try:
                        df_row.reset_index().to_excel(_tmp, index=False)
                        _emb_csv = _gen_emb(_tmp, lm=lm)
                        _emb_df  = _pd.read_csv(_emb_csv, index_col=0)
                        return _emb_df.values.astype('float32')
                    finally:
                        _os.unlink(_tmp)
                        for _ext in ['.emb.csv', f'.{lm}.emb.csv']:
                            _c = _tmp + _ext
                            if _os.path.exists(_c): _os.unlink(_c)
            except Exception:
                _embeddings_fn = None
    else:
        _embeddings_fn = None

    print(f"[Mutagenesis] {model_type.upper()} | {lm} | {_n} antibodies | output → {_mut_dir}/")

    task = getattr(model, 'task', 'classification')
    _saved = []

    for _s in range(_n):
        _row  = data.iloc[_s]
        _bc   = str(data.index[_s])
        _cdr3 = str(_row.get('CDR3', '') or '').upper().replace('-', '')
        if not _cdr3:
            print(f"  [{_s+1}/{_n}] {_bc}: CDR3 missing — skipped")
            continue

        _n_pos = len(_cdr3)
        _n_aa  = len(AMINO_ACIDS)

        _actual = None
        if _label_col:
            try:
                _v = _row.get(_label_col)
                if _v is not None and str(_v) not in ('', 'nan'):
                    _actual = int(float(_v))
            except Exception:
                pass

        print(f"  [{_s+1}/{_n}] {_bc}  CDR3={_cdr3}  ({_n_pos}×{_n_aa} mutants)  "
              + (f"Actual={'PASS' if _actual==1 else 'FAIL'}" if _actual is not None else ""))

        _wt_df = data.iloc[[_s]].copy()
        try:
            _wt_score = float(_predict_single(model, model_type, _wt_df,
                                              _embeddings_fn, task))
        except Exception as _e:
            print(f"    WT score failed: {_e}"); _wt_score = float('nan')

        _mat = np.full((_n_aa, _n_pos), np.nan, dtype=np.float32)
        _vh  = str(_row.get('HSEQ', '') or '')
        _cdr3_start = _vh.find(_cdr3) if _vh else -1

        _first_err = True
        for _pi in range(_n_pos):
            for _ai, _mut_aa in enumerate(AMINO_ACIDS):
                _mcdr3 = _cdr3[:_pi] + _mut_aa + _cdr3[_pi+1:]
                _mrow  = _row.to_dict()
                _mrow['CDR3'] = _mcdr3
                if _cdr3_start >= 0:
                    _mrow['HSEQ'] = (_vh[:_cdr3_start] + _mcdr3 +
                                     _vh[_cdr3_start + _n_pos:])
                _mut_df = pd.DataFrame([_mrow], index=[_bc])
                try:
                    _mat[_ai, _pi] = float(_predict_single(
                        model, model_type, _mut_df, _embeddings_fn, task))
                except Exception as _me:
                    if _first_err:
                        import traceback as _tb
                        print(f"    [Mutagenesis] mutant scoring failed: {_me}")
                        _first_err = False
                    _mat[_ai, _pi] = _wt_score

        _unique_scores = np.unique(_mat[~np.isnan(_mat)])
        if len(_unique_scores) <= 1:
            print(f"    [Mutagenesis] WARNING: all mutant scores = WT ({_wt_score:.4f}) — PLM embedding insensitive")

        _fw = max(9, _n_pos * 0.55 + 3)
        fig, ax = plt.subplots(figsize=(_fw, 7))

        if task == 'classification':
            _norm = _mc.TwoSlopeNorm(vmin=0.0, vcenter=0.5, vmax=1.0)
            _cmap = 'RdBu'
        else:
            _vmin = float(np.nanmin(_mat))
            _vmax = float(np.nanmax(_mat))
            _norm = _mc.TwoSlopeNorm(vmin=_vmin, vcenter=(_vmin+_vmax)/2, vmax=_vmax)
            _cmap = 'coolwarm'

        im = ax.imshow(_mat, cmap=_cmap, norm=_norm, aspect='auto')

        _fsz_cell = max(4.0, min(7.0, 120.0 / max(_n_pos, 1)))
        for _ai in range(_n_aa):
            for _pi in range(_n_pos):
                _v = _mat[_ai, _pi]
                if np.isnan(_v): continue
                _tc = 'white' if (_v < 0.35 or _v > 0.65) else '#333'
                ax.text(_pi, _ai, f"{_v:.2f}", ha='center', va='center',
                        fontsize=_fsz_cell, color=_tc,
                        fontweight='bold' if abs(_v - 0.5) > 0.30 else 'normal')

        for _pi, _wt in enumerate(_cdr3):
            if _wt in AMINO_ACIDS:
                ax.add_patch(plt.Rectangle(
                    (_pi - 0.5, AMINO_ACIDS.index(_wt) - 0.5), 1, 1,
                    fill=False, edgecolor='black', lw=2.0, zorder=5))

        ax.set_xticks(range(_n_pos))
        ax.set_xticklabels([f"{_cdr3[i]}\n{i+1}" for i in range(_n_pos)], fontsize=8.5)
        ax.set_yticks(range(_n_aa))
        ax.set_yticklabels(list(AMINO_ACIDS), fontsize=8)
        ax.set_xlabel('CDR3 position  (WT residue shown above position number)', fontsize=9)
        ax.set_ylabel('Substituted AA', fontsize=9)

        cbar = plt.colorbar(im, ax=ax, shrink=0.85, pad=0.02)
        cbar.set_label('P(PASS)' if task=='classification' else 'Score', fontsize=9, labelpad=6)
        if task == 'classification':
            cbar.set_ticks([0.0, 0.25, 0.5, 0.75, 1.0])
            cbar.set_ticklabels(['0.0\nFAIL', '0.25', '0.50\nborder', '0.75', '1.0\nPASS'])
        cbar.ax.tick_params(labelsize=7.5)

        _score_str = (f"WT P(PASS)={_wt_score:.4f}" if not np.isnan(_wt_score) else "WT score=N/A")
        _act_str   = (f"  |  Actual ({_label_col}) = {'PASS' if _actual==1 else 'FAIL'}"
                      if _actual is not None else "")
        ax.set_title(
            f"DELPHI · CDR3 Mutagenesis Heatmap\n"
            f"ID: {_bc}   {_score_str}{_act_str}\n"
            f"{model_type.upper()} | {lm} | {db_stem}",
            fontsize=9, loc='center', pad=8)
        plt.tight_layout()

        _bc_safe = _bc.replace('/', '_').replace(' ', '_')
        _img_path = os.path.join(_mut_dir, f"{_s+1:04d}_{_bc_safe}_cdr3_mutagenesis.{_fmt}")
        _save_kw = dict(dpi=_dpi, bbox_inches='tight')
        if _fmt == 'tiff':    _save_kw['format'] = 'tiff'
        elif _fmt in ('jpeg','jpg'):
            _save_kw['format'] = 'jpeg'
            _save_kw['pil_kwargs'] = {'quality': 95}
        plt.savefig(_img_path, **_save_kw)
        plt.close()
        _saved.append((_bc, _img_path, _wt_score))
        print(f"    → {os.path.basename(_img_path)}")

    print(f"[Mutagenesis] {len(_saved)} heatmaps → {_mut_dir}/")

    if _make_ppt and _saved:
        try:
            from pptx import Presentation as _Prs
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            from pptx.dml.color import RGBColor

            prs = _Prs()
            prs.slide_width  = Inches(13.33)
            prs.slide_height = Inches(7.5)
            blank = prs.slide_layouts[6]

            for _bc, _img, _wt in _saved:
                slide = prs.slides.add_slide(blank)
                _iw = Inches(11); _ih = Inches(6.5)
                slide.shapes.add_picture(_img,
                    (prs.slide_width - _iw) / 2, Inches(0.4), width=_iw, height=_ih)
                txb = slide.shapes.add_textbox(Inches(0.15), Inches(7.1), Inches(13.0), Inches(0.35))
                tf  = txb.text_frame
                tf.text = (f"{_bc}  |  {model_type.upper()} | {lm} | {db_stem}  |  "
                           f"WT={'%.3f'%_wt if not np.isnan(_wt) else 'N/A'}")
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                p.runs[0].font.size = Pt(7)
                p.runs[0].font.color.rgb = RGBColor(0x88, 0x87, 0x80)

            _ppt = os.path.join(_mut_dir, "cdr3_mutagenesis_all.pptx")
            prs.save(_ppt)
            print(f"[Mutagenesis] PPT ({len(_saved)} slides) → {_ppt}")
        except ImportError:
            print("[Mutagenesis] pip install python-pptx for PPT")
        except Exception as _pe:
            print(f"[Mutagenesis] PPT failed — {_pe}")


def _predict_single(model, model_type, df_row, embeddings_fn, task):
    import numpy as np
    if model_type == 'rf':
        _needs_emb = (hasattr(model, 'fb_') and model.fb_ is not None and
                      model.fb_.feat_cfg.get('embedding', False))
        if _needs_emb and embeddings_fn is not None:
            _emb = embeddings_fn(df_row)
        else:
            _emb = None
        return model.predict_proba(df_row, embeddings=_emb)[0]
    elif model_type == 'xgboost':
        _xgb_row_df = df_row[['HSEQ','LSEQ','CDR3']].copy() if all(
            c in df_row.columns for c in ['HSEQ','LSEQ','CDR3']) else df_row.copy()
        _emb_r = embeddings_fn(df_row) if embeddings_fn else None
        if _emb_r is not None and hasattr(_emb_r, 'reshape'):
            _emb_r = _emb_r.reshape(1, -1)
        if getattr(model, 'task', 'classification') == 'regression':
            return float(model.predict(_xgb_row_df, embeddings=_emb_r)[0])
        return model.predict_proba(_xgb_row_df, embeddings=_emb_r)[0]
    elif model_type in ('transformer_onehot',):
        return model.predict_proba(df_row)[0]
    elif model_type in ('transformer_lm', 'cnn'):
        _emb = embeddings_fn(df_row) if embeddings_fn else None
        if _emb is not None:
            return model.predict_proba(_emb)[0]
        return model.predict_proba(df_row)[0]
    else:
        return model.predict_proba(df_row)[0]


def _find_matching_checkpoints(model_dir: str, target: str, lm: str,
                                model_type: str, db_stem: str,
                                ext: str) -> list:
    import glob, os
    _prefix  = f"FINAL_{target}_{lm}_{model_type}_{db_stem}"
    _pattern = os.path.join(model_dir, f"{_prefix}*{ext}")
    _found   = glob.glob(_pattern)
    _found.sort(key=lambda p: os.path.getmtime(p), reverse=True)
    return _found


def _lookup_registry_path(target, lm, model, db_stem,
                          registry_path="config/model_registry.yaml"):
    """
    Resolve a model checkpoint path from config/model_registry.yaml.

    Matches on target + lm + model (+ db_stem when provided) and returns the
    registered model_path if the file exists on disk. Returns None if there is
    no registry, no matching entry, or the registered file is missing.
    """
    import os as _os
    import yaml as _yaml
    if not _os.path.exists(registry_path):
        return None
    try:
        with open(registry_path) as f:
            reg = _yaml.safe_load(f) or {}
    except Exception:
        return None

    models = reg.get("models", {})
    # Prefer an exact match including db_stem; fall back to target+lm+model
    exact, loose = [], []
    for mid, e in models.items():
        if (e.get("target") == target and e.get("lm") == lm
                and e.get("model") == model):
            mp = e.get("model_path", "")
            if db_stem and db_stem in mid:
                exact.append(mp)
            else:
                loose.append(mp)

    for mp in exact + loose:
        if mp and _os.path.exists(mp):
            print(f"[registry] resolved model → {mp}")
            return mp
    return None


def auto_predict(input_file, target="sec_filter", lm="antiberta2",
                 model_type="xgboost", db_path=None, test_target=None,
                 run_mutagenesis=False, mutagenesis_n=None, threshold=None,
                 model_path=None, **kwargs):
    kwargs['model_path'] = model_path
    print(f"\nPREDICTING: {os.path.basename(input_file)}")
    print(f"Target: {target.upper()} | Model: {model_type.upper()} | LM: {lm}")
    if db_path:
        print(f"Using model trained on: {os.path.basename(db_path)}")

    if input_file.lower().endswith((".xlsx", ".xls")):
        data = pd.read_excel(input_file)
    elif input_file.lower().endswith(".csv"):
        data = pd.read_csv(input_file)
    else:
        raise ValueError(f"Unsupported file format: {input_file}")

    if "BARCODE" not in data.columns:
        data["BARCODE"] = range(len(data))
    data["BARCODE"] = data["BARCODE"].astype(str).str.strip()
    data = data.set_index("BARCODE")

    if "HSEQ" in data.columns:
        data = data.reset_index()
        data = _ensure_cdr3(data, verbose=True)
        data = strip_cdr3_c_prefix(data, verbose=True)   # ← strip C
        data = data.set_index("BARCODE")

    _SEQ_ONLY = {"onehot", "onehot_vh", "onehot_cdr3", "onehot_hcdr3",
                 "biophysical", "kmer", "seq", "none"}

    if lm in _SEQ_ONLY and lm not in ("onehot", "onehot_vh"):
        print(f"[predict] lm='{lm}' — sequence-only mode, no PLM embedding loaded")
        X = data[["HSEQ", "LSEQ", "CDR3"]] if all(
            c in data.columns for c in ["HSEQ", "LSEQ", "CDR3"]) else data
    elif lm not in ["onehot", "onehot_vh"]:
        emb_file = f"{input_file}.{lm}.emb.csv"
        if not os.path.exists(emb_file):
            print("Generating embedding from input file...")
            generate_embedding(input_file, lm=lm)
        print(f"Using embedding: {emb_file}")
        embedding = pd.read_csv(emb_file, index_col=0)
        X, data = _align_embedding(data, embedding, context="predict")
    else:
        required = ["HSEQ", "LSEQ", "CDR3"]
        missing  = [c for c in required if c not in data.columns]
        if missing:
            raise ValueError(f"Missing columns for one-hot: {missing}")

        _nan_like = {"nan", "none", "null", "n/a", "na", "#n/a", "#value!", "#ref!"}
        for _col in required:
            data[_col] = data[_col].fillna("").astype(str)
            _bad = data[_col].str.strip().str.lower().isin(_nan_like)
            if _bad.any():
                print(f"  [fix-nan] {_col}: {_bad.sum():,} nan-like -> ''")
                data.loc[_bad, _col] = ""

        _empty_lseq = data["LSEQ"].str.len() == 0
        if _empty_lseq.any():
            print(f"  [warn] {_empty_lseq.sum():,} antibodies have empty LSEQ — encoded with VL=zeros.")

        _no_hseq = data["HSEQ"].str.strip().str.len() < 5
        if _no_hseq.any():
            print(f"  [fix-nan] Dropping {_no_hseq.sum():,} rows with missing HSEQ")
            data = data[~_no_hseq]
        if len(data) == 0:
            raise ValueError("No valid sequences after cleaning")

        X = data[required]

    _LM_ALIASES = {"onehot_cdr3": "onehot_hcdr3"}
    lm = _LM_ALIASES.get(lm, lm)

    if db_path:
        db_stem = Path(db_path).stem
    elif kwargs.get('model_path'):
        import re as _re2
        _ckpt_stem = Path(kwargs['model_path']).stem
        _m2 = _re2.search(
            r'(?:transformer_lm|transformer_onehot|rf|xgboost|cnn)_(.+?)(?:_ft_|_lora|_plmft|$)',
            _ckpt_stem)
        db_stem = _m2.group(1) if _m2 else _ckpt_stem
    else:
        # Try to extract from model_id if available
        _mid = kwargs.get("model_id") or kwargs.get("model_id_str", "")
        db_stem = (_dbname_from_model_id(_mid, target, lm, model_type)
                   if _mid else "default")

    ext        = ".pt" if model_type in ["cnn", "transformer_onehot", "transformer_lm"] else ".pkl"
    _chain_tag = ""
    if model_type == "rf":
        _SEQ_MODES = {"onehot", "kmer"}
        if lm in _SEQ_MODES or lm in ("biophysical", "none", "seq"):
            try:
                from models.randomforest import RandomForestModel as _RFM_tmp
                _tmp_cfg   = _RFM_tmp("config/random_forest.yaml").config
                _tmp_fb    = __import__('models.randomforest', fromlist=['FeatureBuilder']).FeatureBuilder(_tmp_cfg)
                _chain_tag = _tmp_fb.chain_tag
                if _chain_tag:
                    _chain_tag = f"_{_chain_tag}"
            except Exception:
                _chain_tag = ""

    _explicit_path = kwargs.get('model_path', None)
    if _explicit_path:
        model_path = _explicit_path
        print(f"[load] Using explicit model path: {model_path}")
    else:
        # Resolution order:
        #   1. --model_path (handled above)
        #   2. config/model_registry.yaml  (registered model_path, any directory)
        #   3. {MODEL_DIR}/FINAL_..._{db_stem}{ext}  (filename convention)
        #   4. fuzzy match in MODEL_DIR
        _reg_path = _lookup_registry_path(target, lm, model_type, db_stem)
        if _reg_path:
            model_path = _reg_path
        else:
            _base_path = f"{MODEL_DIR}/FINAL_{target}_{lm}{_chain_tag}_{model_type}_{db_stem}"
            if os.path.exists(f"{_base_path}_regression{ext}"):
                model_path = f"{_base_path}_regression{ext}"
            else:
                model_path = f"{_base_path}{ext}"
            if not os.path.exists(model_path):
                _candidates = _find_matching_checkpoints(
                    MODEL_DIR, target, lm, model_type, db_stem, ext)
                if _candidates:
                    print(f"\n[load] Exact checkpoint not found: {Path(model_path).name}")
                    print(f"[load] Found {len(_candidates)} matching checkpoint(s):")
                    for i, c in enumerate(_candidates):
                        print(f"  [{i}] {Path(c).name}")
                    print(f"[load] Using: [{0}] {Path(_candidates[0]).name}")
                    model_path = _candidates[0]
                else:
                    raise FileNotFoundError(
                        f"Model not found: {model_path}\n"
                        f"Searched in: {MODEL_DIR}\n"
                        f"Also checked config/model_registry.yaml (no usable entry).\n"
                        f"Run --train first, or use --model_path to specify explicitly.")

    if not os.path.exists(model_path):
        raise FileNotFoundError(f"Model not found: {model_path}")

    if model_type == "xgboost":
        model = XGBoostModel.load(model_path)
    elif model_type == "rf":
        model = RandomForestModel.load(model_path)
    elif model_type == "cnn":
        model = CNNModel.load(model_path,
                              embedding_dim=X.shape[1] if lm not in ["onehot","onehot_vh","onehot_hcdr3","onehot_cdr3"] else None)
    elif model_type == "transformer_onehot":
        model = TransformerOneHotModel.load(model_path)
    elif model_type == "transformer_lm":
        model = TransformerLMModel.load(model_path, embedding_dim=X.shape[1])

    if model_type == "rf":
        _SEQ_ONLY_P = {"onehot","onehot_vh","onehot_cdr3","onehot_hcdr3","biophysical","kmer","seq","none"}
        _rf_X_df_p  = data[['HSEQ','LSEQ','CDR3']].copy() if all(
            c in data.columns for c in ['HSEQ','LSEQ','CDR3']) else data.copy()
        _rf_emb_p   = None if lm in _SEQ_ONLY_P else (X.values if hasattr(X, 'values') else X)
        if getattr(model, 'task', 'classification') == 'regression':
            scores = model.predict(_rf_X_df_p, embeddings=_rf_emb_p)
        else:
            scores = model.predict_proba(_rf_X_df_p, embeddings=_rf_emb_p)
    elif model_type == "xgboost":
        _SEQ_ONLY_P = {"onehot","onehot_vh","onehot_cdr3","onehot_hcdr3","biophysical","kmer","seq","none"}
        _xgb_X_df_p = data[['HSEQ','LSEQ','CDR3']].copy() if all(
            c in data.columns for c in ['HSEQ','LSEQ','CDR3']) else data.copy()
        _xgb_emb_p  = None if lm in _SEQ_ONLY_P else (X.values if hasattr(X, 'values') else X)
        if getattr(model, 'task', 'classification') == 'regression':
            scores = model.predict(_xgb_X_df_p, embeddings=_xgb_emb_p)
        else:
            scores = model.predict_proba(_xgb_X_df_p, embeddings=_xgb_emb_p)
    else:
        if getattr(model, 'task', 'classification') == 'regression':
            scores = model.predict(X)
        else:
            scores = model.predict_proba(X)

    _is_regression = getattr(model, 'task', 'classification') == 'regression'
    _reg_tag       = "_regression" if _is_regression else ""
    if _is_regression:
        labels = scores
        print(f"[predict] regression mode — continuous scores")
    elif threshold is not None:
        _thresh = float(threshold)
        labels  = (scores >= _thresh).astype(int)
        print(f"[predict] threshold={_thresh:.4f}  (--threshold override)")
    else:
        _thresh = getattr(model, "recommended_threshold", 0.5)
        labels  = (scores >= _thresh).astype(int)
        print(f"[predict] threshold={_thresh:.4f}")

    data[f"{model_type}_{lm}_{db_stem}_score"] = scores
    data[f"{model_type}_{lm}_{db_stem}_label"] = labels

    if model_type in ("rf", "xgboost") and hasattr(model, 'shap_analysis'):
        try:
            import yaml as _yaml
            _yaml_path = ("config/xgboost.yaml" if model_type == "xgboost"
                          else "config/random_forest.yaml")
            if os.path.exists(_yaml_path):
                with open(_yaml_path) as _yf:
                    _yaml_shap = (_yaml.safe_load(_yf) or {}).get('shap', {})
                if _yaml_shap:
                    model.config['shap'] = _yaml_shap
        except Exception:
            pass
        _shap_cfg = model.config.get('shap', {})
        _shap_ok  = _shap_cfg.get('enabled', True)
        if not _shap_ok:
            print("[SHAP] Skipped — shap.enabled=false in config")
        else:
            try:
                _SEQ_ONLY_SH = {"onehot","onehot_vh","onehot_cdr3","onehot_hcdr3","biophysical","kmer","seq","none"}
                _shap_emb    = None if lm in _SEQ_ONLY_SH else (X.values if hasattr(X, 'values') else X)
                _shap_X_df   = data[['HSEQ','LSEQ','CDR3']].copy() if all(
                    c in data.columns for c in ['HSEQ','LSEQ','CDR3']) else data.copy()
                _shap_top    = _shap_cfg.get('top_features', 30)
                _shap_prefix = f"{Path(input_file).stem}_{target}_{lm}_{model_type}_{db_stem}{_reg_tag}"

                print(f"\n[SHAP] Running on PREDICT set  n={len(_shap_X_df):,}  top={_shap_top}")

                if model.fb_ is None:
                    print("[SHAP] ERROR: FeatureBuilder not found in loaded model.")
                else:
                    if model_type == "xgboost":
                        import models.xgboost as _shap_mod
                    else:
                        import models.randomforest as _shap_mod
                    _orig_mdir = _shap_mod.MODEL_DIR
                    _shap_mod.MODEL_DIR = str(Path(input_file).parent)
                    _shap_actual = None
                    _label_col   = test_target if test_target and test_target in data.columns \
                                   else (target if target in data.columns else None)
                    if _label_col:
                        try:
                            _shap_actual = []
                            for v in data[_label_col].values:
                                if isinstance(v, float) and v != v:
                                    _shap_actual.append(None)
                                else:
                                    try:    _shap_actual.append(int(v))
                                    except: _shap_actual.append(None)
                        except Exception:
                            _shap_actual = None
                    _shap_col = test_target if test_target else target
                    try:
                        model.shap_analysis(
                            _shap_X_df, _shap_emb,
                            output_prefix   = _shap_prefix,
                            split_tag       = "predict",
                            top_n           = _shap_top,
                            barcodes        = list(_shap_X_df.index.astype(str)),
                            actual_labels   = _shap_actual,
                            actual_col_name = _shap_col,
                            lm_name         = lm,
                            db_name         = db_stem,
                        )
                    finally:
                        _shap_mod.MODEL_DIR = _orig_mdir
            except Exception as _se:
                import traceback
                print(f"[SHAP] predict SHAP failed: {_se}")
                print(traceback.format_exc())

    if run_mutagenesis:
        try:
            _run_cdr3_mutagenesis(
                data=data, model=model, model_type=model_type, lm=lm,
                db_stem=db_stem, input_file=input_file, target=target,
                test_target=test_target, embeddings_fn=None, n_override=mutagenesis_n)
        except Exception as _me:
            import traceback
            print(f"[Mutagenesis] failed — {_me}")
            print(traceback.format_exc())

    path        = Path(input_file)
    output_file = path.with_name(
        f"{path.stem}_pred_{target}_{lm}_{model_type}_{db_stem}{_reg_tag}{path.suffix}")
    if path.suffix.lower() in [".xlsx", ".xls"]:
        data.reset_index().to_excel(output_file, index=False)
    else:
        data.reset_index().to_csv(output_file, index=False)

    print(f"Saved predictions to: {output_file}")
    if not _is_regression:
        print(f"Positive rate: {labels.mean():.1%}")

    _eval_label_col  = test_target if (test_target and test_target != target) else target
    _data_with_index = data.reset_index()
    _eval_col_found  = _eval_label_col in _data_with_index.columns
    if not _eval_col_found and target in _data_with_index.columns:
        _eval_label_col = target
        _eval_col_found = True

    if _eval_col_found and not _is_regression:
        try:
            from utils.evaluate_model import evaluate
            _score_col  = f"{model_type}_{lm}_{db_stem}_score"
            _eval_stem  = str(path.with_name(f"{path.stem}_pred_{target}_{lm}_{model_type}_{db_stem}{_reg_tag}"))
            evaluate(file=str(output_file), target=_eval_label_col, score_col=_score_col,
                     cost_fp=1.0, cost_fn=3.0, out=_eval_stem, test_target=test_target,
                     model_type=model_type, lm=lm, db_stem=db_stem,
                     dataset_name=Path(input_file).stem)
        except ImportError:
            print("[eval] utils/evaluate_model.py not found — skipping.")
        except Exception as _e:
            print(f"[eval] WARNING: evaluation failed — {_e}")

    if not _is_regression:
        try:
            from utils.plot_biophysical import plot_biophysical_report
            _bio_stem = str(path.with_name(
                f"{path.stem}_pred_{target}_{lm}_{model_type}_{db_stem}{_reg_tag}"))
            plot_biophysical_report(file=str(output_file), target=target,
                                    test_target=_eval_label_col if _eval_col_found else None,
                                    out=_bio_stem, dataset_name=Path(input_file).stem)
        except ImportError:
            print("[biophys] utils/plot_biophysical.py not found — skipping.")
        except Exception as _e:
            print(f"[biophys] WARNING: {_e}")


_ALL_PLM_LMS = ["ablang", "antiberty", "antiberta2", "antiberta2-cssp", "igbert"]


def auto_predict_multi_lm(input_file, target="psr_filter",
                           lms=None, lm_tag="all",
                           model_type="transformer_lm",
                           db_path=None, test_target=None):
    if lms is None:
        lms = _ALL_PLM_LMS

    print(f"\n{'═'*62}")
    print(f"  MULTI-LM PREDICTION")
    print(f"  File   : {os.path.basename(input_file)}")
    print(f"  Target : {target.upper()}  |  Model: {model_type.upper()}")
    print(f"  LMs    : {lms}")
    if db_path:
        print(f"  DB     : {os.path.basename(db_path)}")
    print(f"{'─'*62}")

    if input_file.lower().endswith((".xlsx", ".xls")):
        base_data = pd.read_excel(input_file)
    elif input_file.lower().endswith(".csv"):
        base_data = pd.read_csv(input_file)
    else:
        raise ValueError(f"Unsupported file format: {input_file}")

    if "BARCODE" not in base_data.columns:
        base_data["BARCODE"] = range(len(base_data))
    base_data["BARCODE"] = base_data["BARCODE"].astype(str).str.strip()

    if "HSEQ" in base_data.columns:
        base_data = _ensure_cdr3(base_data, verbose=True)
        base_data = strip_cdr3_c_prefix(base_data, verbose=True)   # ← strip C

    base_data = base_data.set_index("BARCODE")

    db_stem = (Path(db_path).stem if db_path
             else _dbname_from_model_id(
                 kwargs.get("model_id") or kwargs.get("model_path", "default"),
                 target, "", model_type))
    ext     = ".pt" if model_type in ["cnn", "transformer_onehot", "transformer_lm"] else ".pkl"

    results_summary = []
    failed_lms      = []

    _SEQ_ONLY_MULTI = {"biophysical","kmer","onehot","onehot_vh",
                        "onehot_cdr3","onehot_hcdr3","none","seq"}
    for lm in lms:
        print(f"\n  ── LM: {lm} ──")
        try:
            if lm in _SEQ_ONLY_MULTI:
                _lm_data = base_data.copy()
                X        = _lm_data
            else:
                emb_file = f"{input_file}.{lm}.emb.csv"
                if not os.path.exists(emb_file):
                    print(f"  Generating {lm} embedding ...")
                    generate_embedding(input_file, lm=lm)
                embedding = pd.read_csv(emb_file, index_col=0)
                X, _lm_data = _align_embedding(base_data, embedding, context=f"predict/{lm}")

            _base2     = f"{MODEL_DIR}/FINAL_{target}_{lm}_{model_type}_{db_stem}"
            if os.path.exists(f"{_base2}_regression{ext}"):
                model_path = f"{_base2}_regression{ext}"
            else:
                model_path = f"{_base2}{ext}"
            if not os.path.exists(model_path):
                raise FileNotFoundError(f"Model not found: {model_path}")

            if model_type == "xgboost":
                model = XGBoostModel.load(model_path)
            elif model_type == "rf":
                model = RandomForestModel.load(model_path)
            elif model_type == "cnn":
                model = CNNModel.load(model_path, embedding_dim=X.shape[1])
            elif model_type == "transformer_lm":
                model = TransformerLMModel.load(model_path, embedding_dim=X.shape[1])

            if model_type == "rf":
                _rf_X_df_m = base_data[['HSEQ','LSEQ','CDR3']].copy() if all(
                    c in base_data.columns for c in ['HSEQ','LSEQ','CDR3']) else base_data.copy()
                _rf_X_df_m = _rf_X_df_m.loc[_lm_data.index]
                _SEQ_ONLY_ML = {"biophysical","kmer","onehot","onehot_vh","onehot_cdr3","onehot_hcdr3","none","seq"}
                _rf_emb_m  = None if lm in _SEQ_ONLY_ML else (X.values if hasattr(X, 'values') else X)
                scores = model.predict_proba(_rf_X_df_m, embeddings=_rf_emb_m)
            elif model_type == "xgboost":
                _SEQ_ONLY_M = {"biophysical","kmer","onehot","onehot_vh","onehot_cdr3","onehot_hcdr3","none","seq"}
                _xgb_X_df_m = _lm_data[['HSEQ','LSEQ','CDR3']].copy() if all(
                    c in _lm_data.columns for c in ['HSEQ','LSEQ','CDR3']) else _lm_data.copy()
                _xgb_emb_m  = None if lm in _SEQ_ONLY_M else (X.values if hasattr(X, 'values') else X)
                scores = model.predict_proba(_xgb_X_df_m, embeddings=_xgb_emb_m)
            else:
                scores = model.predict_proba(X)

            _thresh = getattr(model, "recommended_threshold", 0.5)
            labels  = (scores >= _thresh).astype(int)
            pos_rate = labels.mean()
            print(f"  threshold={_thresh:.4f}  n={len(scores):,}  pos_rate={pos_rate:.1%}")

            score_col = f"{model_type}_{lm}_{db_stem}_score"
            label_col = f"{model_type}_{lm}_{db_stem}_label"
            base_data[score_col] = np.nan
            base_data[label_col] = np.nan
            base_data.loc[_lm_data.index, score_col] = scores
            base_data.loc[_lm_data.index, label_col] = labels
            results_summary.append((lm, len(scores), pos_rate, _thresh, _thresh))

        except Exception as e:
            print(f"  [ERROR] {lm} failed: {e}")
            failed_lms.append((lm, str(e)))

    if not results_summary:
        print("\n[multi-lm] ERROR: All LMs failed — no output written.")
        return

    path        = Path(input_file)
    output_file = path.with_name(
        f"{path.stem}_pred_{target}_{lm_tag}_{model_type}_{db_stem}{path.suffix}")
    out_data = base_data.reset_index()
    if path.suffix.lower() in [".xlsx", ".xls"]:
        out_data.to_excel(output_file, index=False)
    else:
        out_data.to_csv(output_file, index=False)

    print(f"\n{'═'*62}")
    print(f"  MULTI-LM SUMMARY")
    print(f"{'─'*62}")
    for row in results_summary:
        lm_name, n, pos, thresh, opt_t = row
        print(f"  {lm_name:25s}  n={n:>7,}  pos={pos:>6.1%}  thresh={thresh:.4f}")
    if failed_lms:
        for lm_name, err in failed_lms:
            print(f"  [FAILED] {lm_name}: {err}")
    print(f"  Output → {output_file}")
    print(f"{'═'*62}\n")


# ===========================================================================
# MAIN
# ===========================================================================

def main():
    parser = argparse.ArgumentParser(description="IPI Antibody Developability Prediction Platform")
    group = parser.add_mutually_exclusive_group(required=True)
    group.add_argument("--predict",         type=str,            help="Predict on new file")
    group.add_argument("--build-embedding", type=str,            help="Generate embeddings only")
    group.add_argument("--kfold",           type=int,            help="Run k-fold CV")
    group.add_argument("--train",           action="store_true", help="Train final model")
    group.add_argument("--split-dataset",   action="store_true", dest="split_dataset",
                       help="Split --db into train + val files.")
    group.add_argument("--list-models",      action="store_true", dest="list_models",
                       help="List models registered in config/model_registry.yaml "
                            "and exit. Optionally filter with --target / --lm / --model.")

    parser.add_argument("--target", type=str, default=None)
    parser.add_argument("--lm", default=None)
    parser.add_argument("--model", default=None,
                        choices=["xgboost","rf","cnn","transformer_onehot","transformer_lm"])
    parser.add_argument("--db",      type=str)
    parser.add_argument("--cluster", type=float, default=0.8, metavar="THRESHOLD")
    parser.add_argument("--split",   type=float, default=0.0, metavar="TRAIN_FRAC")
    parser.add_argument("--val",     type=str,   default=None, metavar="VAL_FILE")
    parser.add_argument("--cluster_col", type=str, default="CDR3",
                        choices=["CDR3", "VH", "VHVL"])
    parser.add_argument("--no-aug",  dest="no_aug", action="store_true", default=True)
    parser.add_argument("--test_target", type=str, default=None)
    parser.add_argument("--mutagenesis", type=int, nargs="?", const=50, default=None, metavar="N")
    parser.add_argument("--threshold", type=float, default=None, metavar="T")
    parser.add_argument("--model_path", type=str, default=None, metavar="PATH")
    parser.add_argument("--finetune", action="store_true", default=False)
    parser.add_argument("--finetune_plm", action="store_true", default=False)
    parser.add_argument("--pretrained", type=str, default=None, metavar="PATH")
    parser.add_argument("--finetune_db", type=str, default=None, metavar="FILE")
    parser.add_argument("--freeze_layers", type=str, default="1", metavar="N")
    parser.add_argument("--freeze_plm_layers", type=int, default=10, metavar="N")
    parser.add_argument("--finetune_lr", type=float, default=1e-6, metavar="LR")
    parser.add_argument("--finetune_epochs", type=int, default=10, metavar="N")
    parser.add_argument("--lr_plm", type=float, default=1e-6, metavar="LR")
    parser.add_argument("--lr_classifier", type=float, default=1e-4, metavar="LR")
    parser.add_argument("--peft", type=str, default="none", choices=["none", "lora"])
    parser.add_argument("--lora_r", type=int, default=8, metavar="R")
    parser.add_argument("--lora_alpha", type=float, default=16.0, metavar="A")
    parser.add_argument("--lora_layers", type=int, nargs="+", default=None, metavar="N")
    parser.add_argument("--cost_fn", type=float, default=3.0, metavar="W")
    parser.add_argument("--cost_fp", type=float, default=1.0, metavar="W")

    args    = parser.parse_args()

    # --lm kmer_vh / kmer_vhvl: aliases selecting the k-mer sequence region (kmer.sequence).
    # Normalize to "kmer" so every downstream --lm gate passes, and stash the region to apply
    # to the model config at the train/kfold dispatch points below (parallels onehot_vh/onehot_cdr3).
    args._kmer_seq = None
    if args.lm and args.lm.lower() in ("kmer_vh", "kmer_vhvl"):
        args._kmer_seq = {"kmer_vh": "VH", "kmer_vhvl": "VHVL"}[args.lm.lower()]
        args.lm = "kmer"

    # ── --list-models: read config/model_registry.yaml and exit ──────────
    if getattr(args, 'list_models', False):
        # Detect which filters the user explicitly passed (vs argparse defaults)
        import sys as _sys
        _argv = " ".join(_sys.argv)
        _f_target = args.target if "--target" in _argv else None
        _f_lm     = args.lm     if "--lm"     in _argv else None
        _f_model  = args.model  if "--model"  in _argv else None
        _list_registered_models(target=_f_target, lm=_f_lm, model=_f_model)
        return

    # ── Infer target/lm/model from a DELPHI-convention --model_path ─────────
    # No silent defaults: if --lm/--model/--target are omitted, we try to read
    # them from the checkpoint filename. Explicit flags always win.
    if getattr(args, 'model_path', None):
        _parsed = _parse_model_filename(args.model_path)
        if _parsed:
            if args.target is None:
                args.target = _parsed['target']
                print(f"[infer] target = {args.target}  (from model filename)")
            if args.lm is None:
                args.lm = _parsed['lm']
                print(f"[infer] lm     = {args.lm}  (from model filename)")
            if args.model is None:
                args.model = _parsed['model']
                print(f"[infer] model  = {args.model}  (from model filename)")

    # ── Require target/lm/model explicitly (no defaults) ───────────────────
    # Only enforced for modes that actually use a model (predict/train/kfold).
    _needs_model_spec = bool(args.predict or args.train or args.kfold)
    if _needs_model_spec:
        _missing = [n for n, v in (("--target", args.target),
                                   ("--lm", args.lm),
                                   ("--model", args.model)) if v is None]
        if _missing:
            parser.error(
                "missing required argument(s): " + ", ".join(_missing) + ".\n"
                "  Pass them explicitly, or use a --model_path whose filename "
                "follows the DELPHI convention\n"
                "  FINAL_{target}_{lm}_{model}_{db_stem}.{pt|pkl} so they can be "
                "inferred.")

    # --build-embedding needs --lm (which PLM to embed with), nothing else.
    if args.build_embedding and args.lm is None:
        parser.error("--build-embedding requires --lm (e.g. --lm ablang).")

    _raw_db = args.db or get_default_db_path()
    if _raw_db and str(_raw_db).lower().endswith(('.pt', '.pkl')):
        if not getattr(args, 'model_path', None):
            args.model_path = _raw_db
            print(f"[db] --db is a checkpoint file → using as --model_path")
        import re as _re
        _stem = Path(_raw_db).stem
        _m = _re.search(
            r'(?:transformer_lm|transformer_onehot|rf|xgboost|cnn)_(.+?)(?:_ft_|_lora|_plmft|$)',
            _stem)
        _inferred_db = _m.group(1) if _m else _stem
        args.db  = None
        db_path  = None
        print(f"[db] Inferred db_stem from checkpoint: '{_inferred_db}'")
    else:
        db_path = _raw_db

    _log_path = _setup_logging(args, db_path)

    _cluster_thresh = args.cluster
    _cluster_col_src = getattr(args, 'cluster_col', 'CDR3').upper()
    _cluster_col_map = {
        'CDR3':  f"HCDR3_CLUSTER_{_cluster_thresh}",
        'HSEQ':  f"VH_CLUSTER_{_cluster_thresh}",
        'VH':    f"VH_CLUSTER_{_cluster_thresh}",
        'VHVL':  f"VHVL_CLUSTER_{_cluster_thresh}",
        'LSEQ':  f"VL_CLUSTER_{_cluster_thresh}",
    }
    _cluster_col = _cluster_col_map.get(_cluster_col_src, f"HCDR3_CLUSTER_{_cluster_thresh}")

    if args.split_dataset:
        if not db_path:
            parser.error("--db required for --split-dataset")
        if not (0.0 < args.split < 1.0):
            parser.error("--split must be in (0,1)")
        split_and_save(db_path=db_path, split=args.split,
                       cluster_thresh=_cluster_thresh,
                       cluster_col=getattr(args, 'cluster_col', 'CDR3'),
                       label_col=args.target)
        return

    if args.build_embedding:
        lms = ["ablang","antiberty","antiberta2","antiberta2-cssp"] if args.lm == "all" else [args.lm]
        if args.lm in ["onehot","onehot_vh"]:
            print("One-hot encoding does not require pre-generation.")
        else:
            for lm in lms:
                generate_embedding(args.build_embedding, lm=lm)
        return

    if args.kfold:
        if not db_path:
            parser.error("--db required for k-fold")
        X, data, y = load_data(db_path, lm=args.lm, label_col=args.target)
        title = f"{args.target.upper()}_{args.model}"

        if _cluster_col not in data.columns:
            if "CDR3" in data.columns:
                print(f"\n[kfold] '{_cluster_col}' not found — computing automatically ...")
                try:
                    from utils.clustering import greedy_clustering_by_levenshtein
                    _seq_col = 'CDR3'
                    data[_cluster_col] = greedy_clustering_by_levenshtein(
                        data[_seq_col].fillna('').tolist(), _cluster_thresh)
                    n_clust = data[_cluster_col].nunique()
                    print(f"[kfold] {n_clust:,} clusters")
                    try:
                        _df_save = read_dataframe(db_path)
                        _df_save[_cluster_col] = data[_cluster_col].values
                        save_dataframe(_df_save, db_path)
                        print(f"[kfold] ✓ Saved {_cluster_col} → {Path(db_path).name}")
                    except Exception as _ce:
                        print(f"[kfold] NOTE: could not save clustering — {_ce}")
                except Exception as _clust_e:
                    print(f"[kfold] WARNING: clustering failed — {_clust_e}")

        if args.model == "xgboost":
            db_stem   = Path(db_path).stem if db_path else ""
            _SEQ_ONLY = {"onehot","onehot_vh","onehot_cdr3","onehot_hcdr3","biophysical","kmer","seq","none"}
            _xgb_emb  = None if args.lm in _SEQ_ONLY else (X.values if hasattr(X, 'values') else X)
            _xgb_X_df = data[['HSEQ','LSEQ','CDR3']].copy() if all(
                c in data.columns for c in ['HSEQ','LSEQ','CDR3']) else data.copy()
            _lm_k = args.lm.lower()
            if _lm_k == "biophysical":
                _xgb_feat = {'embedding': False, 'biophysical': True, 'kmer': False, 'onehot': False}
            elif _lm_k == "kmer":
                _xgb_feat = {'embedding': False, 'biophysical': False, 'kmer': True, 'onehot': False}
                if getattr(args, '_kmer_seq', None): _xgb_feat['_kmer_sequence'] = args._kmer_seq
            elif _lm_k in ("onehot","onehot_vh","onehot_cdr3","onehot_hcdr3"):
                _oh_seq = {"onehot":"VHVL","onehot_vh":"VH","onehot_cdr3":"HCDR3","onehot_hcdr3":"HCDR3"}.get(_lm_k,"VHVL")
                _xgb_feat = {'embedding': False, 'biophysical': False, 'kmer': False, 'onehot': True, '_onehot_sequence': _oh_seq}
            elif _lm_k in ("none","seq"):
                _xgb_feat = {'embedding': False}
            else:
                _xgb_feat = {'embedding': True}
            XGBoostModel.kfold_validation(
                data, _xgb_X_df, y, embeddings=_xgb_emb, embedding_lm=args.lm,
                title=title, kfold=args.kfold, target=args.target,
                cluster_col=_cluster_col, db_stem=db_stem,
                override_features=_xgb_feat,
                cost_fn=getattr(args, 'cost_fn', 3.0), cost_fp=getattr(args, 'cost_fp', 1.0))

        elif args.model == "rf":
            db_stem  = Path(db_path).stem if db_path else ""
            _rf_X_df = data[['HSEQ','LSEQ','CDR3']].copy() if all(
                c in data.columns for c in ['HSEQ','LSEQ','CDR3']) else data.copy()
            _SEQ_ONLY = {"onehot","onehot_vh","onehot_cdr3","onehot_hcdr3","biophysical","kmer","seq","none"}
            _rf_emb = None if args.lm in _SEQ_ONLY else (X.values if hasattr(X, 'values') else X)
            _lm_k = args.lm.lower()
            if _lm_k == "biophysical":
                _rf_feat_override = {'embedding': False, 'biophysical': True, 'kmer': False}
            elif _lm_k == "kmer":
                _rf_feat_override = {'embedding': False, 'biophysical': False, 'kmer': True}
                if getattr(args, '_kmer_seq', None): _rf_feat_override['_kmer_sequence'] = args._kmer_seq
            elif _lm_k in ("none","seq"):
                _rf_feat_override = {'embedding': False}
            elif _lm_k in ("onehot","onehot_vh","onehot_cdr3","onehot_hcdr3"):
                _oh_seq = {"onehot":"VHVL","onehot_vh":"VH","onehot_cdr3":"HCDR3","onehot_hcdr3":"HCDR3"}.get(_lm_k,"VHVL")
                _rf_feat_override = {'embedding': False, 'biophysical': False, 'kmer': False, 'onehot': True, '_onehot_sequence': _oh_seq}
            else:
                _rf_feat_override = {'embedding': True}
            RandomForestModel.kfold_validation(
                data, _rf_X_df, y, embeddings=_rf_emb, embedding_lm=args.lm,
                title=title, kfold=args.kfold, target=args.target,
                cluster_col=_cluster_col, db_stem=db_stem,
                override_features=_rf_feat_override,
                cost_fn=getattr(args, 'cost_fn', 3.0), cost_fp=getattr(args, 'cost_fp', 1.0))

        elif args.model == "transformer_onehot":
            _kf_lm = args.lm if args.lm in ("onehot","onehot_vh") else "onehot"
            TransformerOneHotModel().kfold_validation(
                data, X, y, embedding_lm=_kf_lm, title=title,
                kfold=args.kfold, target=args.target,
                cluster_col=_cluster_col,
                db_stem=Path(db_path).stem if db_path else "")

        elif args.model == "transformer_lm":
            db_stem = Path(db_path).stem if db_path else ""
            TransformerLMModel().kfold_validation(
                db_stem, data, X, y, embedding_lm=args.lm, title=title,
                kfold=args.kfold, target=args.target, cluster_col=_cluster_col)
        return

    if args.train:
        if not db_path:
            parser.error("--db required for training")

        db_stem           = Path(db_path).stem
        _has_explicit_val = bool(getattr(args, 'val', None))
        _has_split        = 0.0 < args.split < 1.0
        val_X = val_y = val_db_path = data_val = None

        if _has_explicit_val:
            X, data, y             = load_data(db_path,  lm=args.lm, label_col=args.target)
            X_val, data_val, y_val = load_data(args.val, lm=args.lm, label_col=args.target)
            val_X, val_y, val_db_path = X_val, y_val, args.val
        elif _has_split:
            train_path, val_path, _, _ = split_and_save(
                db_path=db_path, split=args.split,
                cluster_thresh=_cluster_thresh,
                cluster_col=getattr(args, 'cluster_col', 'CDR3'),
                label_col=args.target)
            X, data, y             = load_data(train_path, lm=args.lm, label_col=args.target)
            X_val, data_val, y_val = load_data(val_path,   lm=args.lm, label_col=args.target)
            val_X, val_y, val_db_path = X_val, y_val, val_path
        else:
            X, data, y = load_data(db_path, lm=args.lm, label_col=args.target)

        if args.model == "xgboost":
            model = XGBoostModel(verbose=False)
            _lm_xgb = args.lm.lower()
            if _lm_xgb == "biophysical":
                model.config['features'].update({'embedding': False, 'biophysical': True, 'kmer': False, 'onehot': False})
            elif _lm_xgb == "kmer":
                model.config['features'].update({'embedding': False, 'biophysical': False, 'kmer': True, 'onehot': False})
                if getattr(args, '_kmer_seq', None): model.config.setdefault('kmer', {})['sequence'] = args._kmer_seq
            elif _lm_xgb in ("onehot","onehot_vh","onehot_cdr3","onehot_hcdr3"):
                _oh_seq_xgb = {"onehot":"VHVL","onehot_vh":"VH","onehot_cdr3":"HCDR3","onehot_hcdr3":"HCDR3"}.get(_lm_xgb,"VHVL")
                model.config['features'].update({'embedding': False, 'biophysical': False, 'kmer': False, 'onehot': True})
                model.config.setdefault('onehot', {})['sequence'] = _oh_seq_xgb
            elif _lm_xgb in ("none","seq"):
                model.config['features']['embedding'] = False
            else:
                model.config['features']['embedding'] = True
            XGBoostModel.print_config_report(model.config)

        elif args.model == "rf":
            model = RandomForestModel(verbose=False)
            _lm = args.lm.lower()
            if _lm == "biophysical":
                model.config['features'].update({'embedding': False, 'biophysical': True, 'kmer': False})
            elif _lm == "kmer":
                model.config['features'].update({'embedding': False, 'biophysical': False, 'kmer': True})
                if getattr(args, '_kmer_seq', None): model.config.setdefault('kmer', {})['sequence'] = args._kmer_seq
            elif _lm in ("none","seq"):
                model.config['features']['embedding'] = False
            elif _lm in ("onehot","onehot_vh","onehot_cdr3","onehot_hcdr3"):
                _oh_seq = {"onehot":"VHVL","onehot_vh":"VH","onehot_cdr3":"HCDR3","onehot_hcdr3":"HCDR3"}.get(_lm,"VHVL")
                model.config.setdefault('features', {}).update({'embedding': False, 'biophysical': False, 'kmer': False, 'onehot': True})
                model.config.setdefault('onehot', {})['sequence'] = _oh_seq
            else:
                model.config['features']['embedding'] = True
            from models.randomforest import RandomForestModel as _RFM
            _RFM.print_config_report(model.config)

        elif args.model == "cnn":
            model = CNNModel()
        elif args.model == "transformer_onehot":
            model = TransformerOneHotModel()
            _t_lm = args.lm if args.lm in ("onehot","onehot_vh") else "onehot"
            model.set_lm_mode(_t_lm)
        elif args.model == "transformer_lm":
            model = TransformerLMModel()

        _vkw = {}
        if val_X is not None and args.model in ("transformer_lm","transformer_onehot","cnn"):
            _vkw = {"val_X": val_X, "val_y": val_y}

        if args.model == "transformer_onehot":
            model.train(X, y, target=args.target, db_stem=db_stem,
                        cluster_col=_cluster_col, no_aug=args.no_aug, **_vkw)
        elif args.model == "transformer_lm":
            model.train(X, y, target=args.target, db_stem=db_stem,
                        embedding_lm=args.lm, cluster_col=_cluster_col, **_vkw)
        elif args.model == "rf":
            _rf_train_X_df = data[['HSEQ','LSEQ','CDR3']].copy() if all(
                c in data.columns for c in ['HSEQ','LSEQ','CDR3']) else data.copy()
            _SEQ_ONLY_TRAIN = {"biophysical","kmer","onehot","onehot_vh","onehot_cdr3","onehot_hcdr3","none","seq"}
            _rf_train_emb   = None if args.lm in _SEQ_ONLY_TRAIN else (X.values if hasattr(X, 'values') else X)
            _rf_val_X_df = _rf_val_emb = None
            if val_X is not None and val_y is not None and data_val is not None:
                _rf_val_X_df = data_val[['HSEQ','LSEQ','CDR3']].copy() if all(
                    c in data_val.columns for c in ['HSEQ','LSEQ','CDR3']) else data_val.copy()
                _rf_val_emb  = None if args.lm in _SEQ_ONLY_TRAIN else (val_X.values if hasattr(val_X, 'values') else val_X)
            _y_uniq_rf   = len(set(y.tolist() if hasattr(y, 'tolist') else list(y)))
            _rf_task_tag = "_regression" if _y_uniq_rf > 2 else ""
            _rf_full_stem = f"{args.target}_{args.lm}_rf_{db_stem}{_rf_task_tag}"
            model.train(_rf_train_X_df, y,
                        embeddings=_rf_train_emb, val_X=_rf_val_X_df,
                        val_y=val_y if val_y is not None else None,
                        val_embeddings=_rf_val_emb,
                        target=_rf_full_stem, target_col=args.target, embedding_lm=args.lm)
        elif args.model == "xgboost":
            _xgb_train_X_df = data[['HSEQ','LSEQ','CDR3']].copy() if all(
                c in data.columns for c in ['HSEQ','LSEQ','CDR3']) else data.copy()
            _SEQ_ONLY_XGB = {"biophysical","kmer","onehot","onehot_vh","onehot_cdr3","onehot_hcdr3","none","seq"}
            _xgb_train_emb = None if args.lm in _SEQ_ONLY_XGB else (X.values if hasattr(X, 'values') else X)
            _xgb_val_X_df = _xgb_val_emb = None
            if val_X is not None and val_y is not None and data_val is not None:
                _xgb_val_X_df = data_val[['HSEQ','LSEQ','CDR3']].copy() if all(
                    c in data_val.columns for c in ['HSEQ','LSEQ','CDR3']) else data_val.copy()
                _xgb_val_emb = None if args.lm in _SEQ_ONLY_XGB else (val_X.values if hasattr(val_X, 'values') else val_X)
            _y_uniq_xgb    = len(set(y.tolist() if hasattr(y, 'tolist') else list(y)))
            _xgb_task_tag  = "_regression" if _y_uniq_xgb > 2 else ""
            _xgb_full_stem = f"{args.target}_{args.lm}_xgboost_{db_stem}{_xgb_task_tag}"
            model.train(_xgb_train_X_df, y,
                        embeddings=_xgb_train_emb, val_X=_xgb_val_X_df,
                        val_y=val_y if val_y is not None else None,
                        val_embeddings=_xgb_val_emb,
                        target=_xgb_full_stem, target_col=args.target, embedding_lm=args.lm)
        else:
            model.train(X, y, embedding_lm=args.lm)

        ext  = ".pt" if args.model in ["cnn","transformer_onehot","transformer_lm"] else ".pkl"
        _task_suffix = ""
        if hasattr(model, 'task') and model.task == 'regression':
            _task_suffix = "_regression"
        path = f"{MODEL_DIR}/FINAL_{args.target}_{args.lm}_{args.model}_{db_stem}{_task_suffix}{ext}"
        model.save(path)
        print(f"FINAL MODEL SAVED: {path}")

        # ── Register in model_registry.yaml ──────────────────────────────
        _reg_thresh = getattr(model, 'recommended_threshold', None)
        _register_model(
            model_path=path,
            target=args.target,
            lm=args.lm,
            model=args.model,
            trainset=(args.db or ""),
            threshold=(_reg_thresh if _reg_thresh not in (None, 0.5) else None),
            epochs=getattr(args, 'finetune_epochs', None) if False else None,
        )

    if args.predict:
        _VALID_PLM_LMS = {"ablang","antiberty","antiberta2","antiberta2-cssp","igbert"}
        _SEQ_ONLY_LMS  = {"onehot","onehot_vh","onehot_cdr3","onehot_hcdr3","biophysical","kmer","seq","none"}
        _VALID_ALL_LMS = _VALID_PLM_LMS | _SEQ_ONLY_LMS

        _lm_raw = args.lm.strip()
        if _lm_raw == "all":
            _lm_list = _ALL_PLM_LMS
        elif "," in _lm_raw:
            _lm_list = [x.strip() for x in _lm_raw.split(",") if x.strip()]
        else:
            _lm_list = None

        if _lm_list is not None:
            _lm_tag = "all" if set(_lm_list) == set(_ALL_PLM_LMS) else "_".join(_lm_list)
            auto_predict_multi_lm(args.predict, target=args.target, lms=_lm_list,
                                   lm_tag=_lm_tag, model_type=args.model,
                                   db_path=args.db, test_target=args.test_target)
        else:
            auto_predict(args.predict, target=args.target, lm=_lm_raw,
                         model_type=args.model, db_path=args.db,
                         test_target=args.test_target,
                         run_mutagenesis=args.mutagenesis is not None,
                         mutagenesis_n=args.mutagenesis,
                         threshold=args.threshold,
                         model_path=getattr(args, 'model_path', None))


if __name__ == "__main__":
    try:
        main()
    finally:
        if isinstance(sys.stdout, _Tee):
            ts_end = datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S')
            print(f"\n[log] Finished: {ts_end}")
            sys.stdout.flush()
            sys.stdout.close()
            sys.stdout = sys.__stdout__
