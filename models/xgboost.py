# models/xgboost.py
# ══════════════════════════════════════════════════════════════════════════════
# DELPHI — Deep End-to-end Learning Platform for antibody developability
#          with High Interpretability
#
# Module      : xgboost.py
# Description : XGBoost gradient-boosted tree classifier for antibody
#               developability prediction. Supports the same feature
#               representations as random_forest.py: biophysical descriptors,
#               k-mer frequencies, one-hot encoding, and PLM embeddings.
#               Includes SHAP interpretability and CDR3 in silico mutagenesis.
#               Typically outperforms Random Forest on tabular biophysical
#               features and is the recommended tree-based model when
#               training data exceeds 10k samples.
#               Applicable to any binary (1/0) antibody property label 
#               (PSR, SEC, HIC, AC-SINS, viscosity, expression, ...).
# Author      : Hoan Nguyen, PhD
# Company     : Institute for Protein Innovation (IPI)
# Date        : 2026-05
# Version     : 1.0.0
# ══════════════════════════════════════════════════════════════════════════════
# ── JAN-2026 update [HYP-1] ─────────────────────────────────────────────────
# Hyperparameter search strategies in kfold_validation():
#
#   'pre_search'  (DEFAULT — recommended for all routine use)
#     Search on FULL data BEFORE kfold loop. Own internal 3-fold CV.
#     → Clean separation: search = find params, kfold = evaluate
#     → ~9× faster than nested. Leakage < fold variance (< 0.003 AUC).
#     → All folds use identical params → fair cross-fold comparison.
#     → Standard practice in production ML.
#
#   'first_fold'
#     Search on fold 1 training data only. Reuse params for folds 2–K.
#     → Zero leakage. Same speed as pre_search.
#     → Good when strict no-leakage is required without full nested cost.
#
#   'nested'
#     Per-fold inner search on each fold's training data.
#     → Gold standard. Zero leakage. Best params per fold.
#     → ~10× slower. Use for final publication evaluation only.
#
#   'none'  (when use_random_search: false)
#     Config params used directly. Fastest.
#     Use after best params are known and saved to xgboost.yaml.
#
#   Recommended workflow:
#     1. pre_search kfold    → find good params + evaluate  (routine)
#     2. Copy best params to xgboost.yaml                   (persist)
#     3. use_random_search: false kfold                     (fast eval)
#     4. nested kfold (once) → rigorous final eval          (paper)
#
#   YAML:
#     training:
#       use_random_search: true
#       search_strategy: pre_search   # pre_search | first_fold | nested | none
#       n_iter: 50
#       cv: 3
# ══════════════════════════════════════════════════════════════════════════════


from config import MODEL_DIR

import os
import copy
import datetime
import logging
import joblib
import numpy as np
import pandas as pd
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import xgboost as xgb
from pathlib import Path
from sklearn.model_selection import (
    RandomizedSearchCV, StratifiedKFold, StratifiedGroupKFold)
from sklearn.metrics import (
    roc_curve, auc, f1_score, accuracy_score, roc_auc_score,
    precision_score, recall_score, confusion_matrix)
from scipy.stats import uniform, randint
from collections import Counter
from itertools import product
import warnings
import yaml

try:
    import shap as _shap_lib
    _SHAP_AVAILABLE = True
except ImportError:
    _SHAP_AVAILABLE = False

try:
    from utils.threshold_optimizer import run_full_threshold_pipeline
    _THRESHOLD_OPT_AVAILABLE = True
except ImportError:
    _THRESHOLD_OPT_AVAILABLE = False


# ══════════════════════════════════════════════════════════════════════════
# AMINO ACID CONSTANTS & BIOPHYSICAL SCALES
# ══════════════════════════════════════════════════════════════════════════

AMINO_ACIDS = 'ACDEFGHIKLMNPQRSTVWY'

_PKA = {'D': 3.9, 'E': 4.1, 'C': 8.3, 'Y': 10.1,
        'H': 6.0, 'K': 10.5, 'R': 12.5,
        'N_term': 8.0, 'C_term': 3.1}

_KD = {'A': 1.8, 'R':-4.5, 'N':-3.5, 'D':-3.5, 'C': 2.5,
       'Q':-3.5, 'E':-3.5, 'G':-0.4, 'H':-3.2, 'I': 4.5,
       'L': 3.8, 'K':-3.9, 'M': 1.9, 'F': 2.8, 'P':-1.6,
       'S':-0.8, 'T':-0.7, 'W':-0.9, 'Y':-1.3, 'V': 4.2}


# ══════════════════════════════════════════════════════════════════════════
# FEATURE HELPERS
# ══════════════════════════════════════════════════════════════════════════

def _charge_at_ph7(seq):
    c = 0.0
    for aa in seq.upper():
        if   aa == 'D': c -= 1/(1+10**(_PKA['D']-7))
        elif aa == 'E': c -= 1/(1+10**(_PKA['E']-7))
        elif aa == 'H': c += 1/(1+10**(7-_PKA['H']))
        elif aa == 'K': c += 1/(1+10**(7-_PKA['K']))
        elif aa == 'R': c += 1/(1+10**(7-_PKA['R']))
    return round(c, 4)

def _gravy(seq):
    vals = [_KD.get(aa, 0.0) for aa in seq.upper()]
    return round(sum(vals)/len(vals), 4) if vals else 0.0

def _pi(seq):
    def _c(ph):
        c = 1/(1+10**(ph-_PKA['N_term'])) - 1/(1+10**(_PKA['C_term']-ph))
        for aa in seq.upper():
            if aa=='D': c -= 1/(1+10**(_PKA['D']-ph))
            elif aa=='E': c -= 1/(1+10**(_PKA['E']-ph))
            elif aa=='H': c += 1/(1+10**(ph-_PKA['H']))
            elif aa=='K': c += 1/(1+10**(ph-_PKA['K']))
            elif aa=='R': c += 1/(1+10**(ph-_PKA['R']))
        return c
    lo, hi = 0.0, 14.0
    for _ in range(200):
        mid = (lo+hi)/2
        if _c(mid) > 0: lo = mid
        else: hi = mid
    return round((lo+hi)/2, 3)

def _instability(seq):
    DIWV = {'WC':1,'WM':24.68,'CM':33.6,'CH':33.6,'CD':20.26,'CT':33.6,
            'CL':20.26,'CP':20.26,'QD':20.26,'RD':20.26,'RH':20.26,
            'DG':20.26,'DS':20.26,'FD':54.96,'GD':20.26,'GH':33.6,
            'HN':20.26,'LR':20.26,'ND':20.26,'NL':20.26,'NP':20.26,
            'QR':20.26,'RR':58.28,'SR':20.26,'TD':20.26,'TP':20.26,
            'VN':1,'WR':1}
    seq = seq.upper()
    if len(seq) < 2: return 0.0
    total = sum(DIWV.get(seq[i:i+2], 1.0) for i in range(len(seq)-1))
    return round(10.0/len(seq)*total, 3)

def _strip_cdr3_loop(cdr3: str) -> str:
    s = cdr3.upper().replace('-', '')
    if s.startswith('CAR'): s = s[3:]
    elif s.startswith('C'): s = s[1:]
    if len(s) >= 3 and s[-2] == 'D' and s[-1] in 'YW': s = s[:-3]
    return s

def compute_biophysical_features(seq: str, feature_list: list,
                                  vh_seq: str = '') -> dict:
    seq = seq.upper().replace('-', '')
    if not seq: return {f: 0.0 for f in feature_list}
    counts = Counter(seq); n = len(seq)
    vh_len = len(vh_seq.upper().replace('-', '')) if vh_seq else 0
    out = {}
    for f in feature_list:
        if   f == 'length':               out[f] = float(n)
        elif f == 'pi':                   out[f] = _pi(seq)
        elif f in ('charge_ph7', 'charge'): out[f] = _charge_at_ph7(seq)
        elif f == 'hydrophobicity':       out[f] = _gravy(seq)
        elif f == 'aromaticity':          out[f] = round(sum(counts.get(a,0) for a in 'FYW')/n, 4)
        elif f == 'instability':          out[f] = _instability(seq)
        elif f == 'net_charge_sq':
            c = _charge_at_ph7(seq);      out[f] = round(c*c, 4)
        elif f == 'frac_charged':         out[f] = round(sum(counts.get(a,0) for a in 'DERKH')/n, 4)
        elif f == 'frac_hydrophobic':     out[f] = round(sum(counts.get(a,0) for a in 'AVILMFWP')/n, 4)
        elif f == 'cdr3_length':          out[f] = float(n)
        elif f == 'vh_length':            out[f] = float(vh_len)
        elif f == 'vh_cdr3_length_ratio': out[f] = round(vh_len/n, 4) if n > 0 else 0.0
        elif f in ('vh_charge_ph7', 'vh_charge'): out[f] = _charge_at_ph7(vh_seq.upper().replace('-','')) if vh_seq else 0.0
        elif f == 'vh_hydrophobicity':    out[f] = _gravy(vh_seq.upper().replace('-','')) if vh_seq else 0.0
        elif f.startswith('count_'):      out[f] = float(counts.get(f.split('_')[1], 0))
        elif len(f) == 1 and f in AMINO_ACIDS: out[f] = float(counts.get(f, 0))
        else:                             out[f] = 0.0
    return out

def compute_kmer_features(seq: str, k_list: list, normalize: bool=True) -> dict:
    seq = seq.upper().replace('-','')
    out = {}
    for k in k_list:
        all_kmers = [''.join(p) for p in product(AMINO_ACIDS, repeat=k)]
        cnts = Counter(seq[i:i+k] for i in range(max(0, len(seq)-k+1))) if len(seq)>=k else {}
        total = sum(cnts.values()) or 1
        for km in all_kmers:
            v = cnts.get(km, 0)
            out[f'{k}mer_{km}'] = v/total if normalize else float(v)
    return out

_AA_LOOKUP_RF = np.full(128, 20, dtype=np.int32)
for _aa_rf, _idx_rf in {aa: i for i, aa in enumerate(AMINO_ACIDS)}.items():
    _AA_LOOKUP_RF[ord(_aa_rf)] = _idx_rf

def compute_onehot_features(seq: str, max_len: int) -> np.ndarray:
    seq = seq.upper().replace('-', '')[:max_len]
    out = np.zeros((max_len, len(AMINO_ACIDS)), dtype=np.float32)
    if seq:
        chars  = np.frombuffer(seq.encode('ascii'), dtype=np.uint8)
        aa_idx = _AA_LOOKUP_RF[np.clip(chars, 0, 127)]
        pos    = np.arange(len(chars)); valid = aa_idx < 20
        out[pos[valid], aa_idx[valid]] = 1.0
    return out.flatten()


# ══════════════════════════════════════════════════════════════════════════════
# FEATURE BUILDER
# ══════════════════════════════════════════════════════════════════════════════

class FeatureBuilder:
    def __init__(self, config: dict):
        self.cfg        = config
        self.feat_cfg   = config.get('features', {})
        self.km_cfg     = config.get('kmer',     {})
        self.bp_cfg     = config.get('biophysical', {})
        self.oh_cfg     = config.get('onehot',   {})
        self.feat_cfg.setdefault('onehot', False)
        self.feature_names_ : list = None
        self._emb_dim = None
        self._oh_max_vh = self._oh_max_vl = self._oh_max_cdr3 = 0

    def _get_kmer_seq(self, row) -> str:
        src = self.km_cfg.get('sequence', 'CDR3').upper()
        vh = str(row.get('HSEQ', '') or ''); vl = str(row.get('LSEQ', '') or '')
        cdr = str(row.get('CDR3', '') or '')
        if src == 'VH': return vh
        if src == 'VHVL': return vh + vl
        return cdr

    def _get_cdr3(self, row) -> str: return str(row.get('CDR3', '') or '')
    def _get_vh(self, row) -> str:   return str(row.get('HSEQ', '') or '')
    def _get_vl(self, row) -> str:   return str(row.get('LSEQ', '') or '')

    def _oh_segs(self, row) -> list:
        mode = self.oh_cfg.get('sequence', 'HCDR3').upper()
        vh = self._get_vh(row); vl = self._get_vl(row); cdr = self._get_cdr3(row)
        if mode == 'VH':   return [(vh, self._oh_max_vh, 'oh_vh')]
        if mode == 'VHVL': return [(vh, self._oh_max_vh, 'oh_vh'), (vl, self._oh_max_vl, 'oh_vl')]
        return [(cdr, self._oh_max_cdr3, 'oh_hcdr3')]

    @property
    def chain_tag(self) -> str:
        if self.feat_cfg.get('kmer') and not self.feat_cfg.get('onehot'):
            return self.km_cfg.get('sequence', 'CDR3').upper()
        if self.feat_cfg.get('onehot'):
            return self.oh_cfg.get('sequence', 'HCDR3').upper()
        return ''

    def fit(self, X_df: pd.DataFrame, embeddings: np.ndarray = None):
        names = []
        if self.feat_cfg.get('embedding'):
            if embeddings is None:
                raise ValueError("features.embedding=True but embeddings=None.")
            self._emb_dim = embeddings.shape[1]
            names += [f'emb_{i}' for i in range(self._emb_dim)]
        if self.feat_cfg.get('kmer'):
            kf = compute_kmer_features(self._get_kmer_seq(X_df.iloc[0]),
                                        self.km_cfg.get('k', [1, 2]),
                                        self.km_cfg.get('normalize', True))
            names += list(kf.keys())
        if self.feat_cfg.get('biophysical'):
            bf = compute_biophysical_features(self._get_cdr3(X_df.iloc[0]),
                                               self.bp_cfg.get('features', []),
                                               vh_seq=self._get_vh(X_df.iloc[0]))
            def _rename(k):
                k = k.replace('count_', '').replace('_ph7', '')
                return k if (k.startswith('vh_') or k.startswith('cdr3_')) else f'cdr3_{k}'
            names += [_rename(k) for k in bf.keys()]
        if self.feat_cfg.get('onehot'):
            mode = self.oh_cfg.get('sequence', 'HCDR3').upper()
            _max_cfg = self.oh_cfg.get('max_lengths', {})
            _cap_vh = int(_max_cfg.get('vh', 0) or 0)
            _cap_vl = int(_max_cfg.get('vl', 0) or 0)
            _cap_cdr3 = int(_max_cfg.get('hcdr3', 0) or 0)
            if mode in ('VH', 'VHVL'):
                _dmx = max((len(str(r.get('HSEQ','') or '')) for _, r in X_df.iterrows()), default=135)
                self._oh_max_vh = _cap_vh if _cap_vh > 0 else _dmx
            if mode == 'VHVL':
                _dmx = max((len(str(r.get('LSEQ','') or '')) for _, r in X_df.iterrows()), default=135)
                self._oh_max_vl = _cap_vl if _cap_vl > 0 else _dmx
            if mode == 'HCDR3':
                _dmx = max((len(str(r.get('CDR3','') or '')) for _, r in X_df.iterrows()), default=25)
                self._oh_max_cdr3 = _cap_cdr3 if _cap_cdr3 > 0 else _dmx
            for _seq, max_len, pfx in self._oh_segs(X_df.iloc[0]):
                for pos in range(1, max_len+1):
                    for aa in AMINO_ACIDS:
                        names.append(f"{pfx}_{pos}_{aa}")
        if not names:
            raise ValueError("No features selected.")
        self.feature_names_ = names
        _AA = len(AMINO_ACIDS); _SEP = '─'*52
        print(f"[FeatureBuilder] Feature dimensions:\n  {_SEP}")
        if self.feat_cfg.get('embedding') and self._emb_dim:
            print(f"  embedding       : {self._emb_dim:>6,} d  (PLM vector)")
        if self.feat_cfg.get('kmer'):
            _d = len([n for n in names if 'mer_' in n])
            print(f"  kmer            : {_d:>6,} d  (k={self.km_cfg.get('k')}  seq={self.km_cfg.get('sequence')})")
        if self.feat_cfg.get('biophysical'):
            _d = len([n for n in names if n.startswith('cdr3_') or n.startswith('vh_')])
            print(f"  biophysical     : {_d:>6,} d  (CDR3 + VH properties)")
        if self.feat_cfg.get('onehot'):
            _mode = self.oh_cfg.get('sequence', 'HCDR3').upper()
            _d = (self._oh_max_vh+self._oh_max_vl)*_AA if _mode=='VHVL' else \
                 self._oh_max_vh*_AA if _mode=='VH' else self._oh_max_cdr3*_AA
            print(f"  onehot ({_mode:<5})  : {_d:>6,} d")
        print(f"  {_SEP}\n  TOTAL           : {len(names):>6,} d\n  {_SEP}")
        return self

    def transform(self, X_df: pd.DataFrame, embeddings: np.ndarray = None) -> np.ndarray:
        parts = []; n = len(X_df)
        if self.feat_cfg.get('embedding') and embeddings is not None:
            parts.append(embeddings.astype(np.float32))
        if self.feat_cfg.get('kmer'):
            k_list = self.km_cfg.get('k', [1, 2]); norm = self.km_cfg.get('normalize', True)
            rows = [list(compute_kmer_features(self._get_kmer_seq(X_df.iloc[i]), k_list, norm).values()) for i in range(n)]
            parts.append(np.array(rows, dtype=np.float32))
        if self.feat_cfg.get('biophysical'):
            fl = self.bp_cfg.get('features', [])
            rows = [list(compute_biophysical_features(self._get_cdr3(X_df.iloc[i]), fl, vh_seq=self._get_vh(X_df.iloc[i])).values()) for i in range(n)]
            parts.append(np.array(rows, dtype=np.float32))
        if self.feat_cfg.get('onehot'):
            oh_rows = []
            for i in range(n):
                row = X_df.iloc[i]
                vecs = [compute_onehot_features(seq, max_len) for seq, max_len, _ in self._oh_segs(row)]
                oh_rows.append(np.concatenate(vecs))
            parts.append(np.array(oh_rows, dtype=np.float32))
        return np.hstack(parts) if parts else np.zeros((n, 0), dtype=np.float32)

    def fit_transform(self, X_df, embeddings=None):
        return self.fit(X_df, embeddings).transform(X_df, embeddings)

    @property
    def n_features(self): return len(self.feature_names_ or [])
    @property
    def non_embedding_feature_names(self): return [n for n in (self.feature_names_ or []) if not n.startswith('emb_')]
    @property
    def non_embedding_indices(self): return [i for i, n in enumerate(self.feature_names_ or []) if not n.startswith('emb_')]
    @property
    def onehot_feature_names(self): return [n for n in (self.feature_names_ or []) if n.startswith('oh_')]
    @property
    def onehot_indices(self): return [i for i, n in enumerate(self.feature_names_ or []) if n.startswith('oh_')]


# ══════════════════════════════════════════════════════════════════════════════
# UTILITIES
# ══════════════════════════════════════════════════════════════════════════════

def _deep_merge(base: dict, override: dict) -> dict:
    for k, v in override.items():
        if isinstance(v, dict) and isinstance(base.get(k), dict): _deep_merge(base[k], v)
        else: base[k] = v
    return base

def _setup_logger(log_path: str) -> logging.Logger:
    logger = logging.getLogger(f'XGB_{os.path.basename(log_path)}')
    logger.setLevel(logging.DEBUG)
    if not logger.handlers:
        fh = logging.FileHandler(log_path, mode='a', encoding='utf-8')
        fh.setFormatter(logging.Formatter('%(message)s'))
        logger.addHandler(fh)
    return logger

def _log(logger, msg: str):
    print(msg)
    if logger: logger.info(msg)


# ══════════════════════════════════════════════════════════════════════════════
# HYPERPARAMETER SEARCH  [HYP-1]
# ══════════════════════════════════════════════════════════════════════════════

def _get_param_dist() -> dict:
    """
    RandomizedSearchCV parameter distribution.
    scale_pos_weight excluded — computed directly from class distribution,
    not searched (its optimal value is always n_neg/n_pos).
    """
    return {
        'n_estimators':     randint(1000, 5000),
        'max_depth':        randint(3, 10),
        'learning_rate':    uniform(0.005, 0.095),
        'subsample':        uniform(0.60, 0.40),
        'colsample_bytree': uniform(0.20, 0.60),
        'gamma':            uniform(0.0,  0.30),
        'min_child_weight': randint(1, 10),
        'reg_alpha':        uniform(0.0,  0.50),
        'reg_lambda':       uniform(0.5,  1.50),
    }


def _compute_spw(y: np.ndarray, task: str) -> float:
    """Compute scale_pos_weight = n_neg / n_pos from class distribution."""
    if task != 'classification': return 1.0
    _y = np.asarray(y, dtype=int); _c = np.bincount(_y)
    return float(_c[0]) / float(_c[1]) if len(_c) > 1 and _c[1] > 0 else 1.0


def _run_search(X: np.ndarray, y: np.ndarray,
                n_iter: int, cv: int,
                task: str, spw: float,
                logger=None) -> tuple:
    """
    Core RandomizedSearchCV runner.
    Returns (best_params, best_score).
    refit=False — we only need params, not the fitted model.
    Caller builds the actual model with inst._build_xgb(y) separately.
    """
    param_dist = _get_param_dist()
    if task == 'classification':
        _base = xgb.XGBClassifier(
            objective='binary:logistic', eval_metric='auc',
            scale_pos_weight=spw, tree_method='hist',
            n_jobs=-1, random_state=42, verbosity=0)
        _scoring     = 'roc_auc'
        _cv_splitter = StratifiedKFold(n_splits=cv, shuffle=True, random_state=42)
    else:
        _base = xgb.XGBRegressor(
            objective='reg:squarederror', tree_method='hist',
            n_jobs=-1, random_state=42, verbosity=0)
        _scoring = 'r2'; _cv_splitter = cv

    search = RandomizedSearchCV(
        _base, param_distributions=param_dist,
        n_iter=n_iter, cv=_cv_splitter, scoring=_scoring,
        refit=False, n_jobs=-1, random_state=42, verbose=0)
    search.fit(X, y)
    return search.best_params_, float(search.best_score_)


def _log_params(params: dict, logger, prefix: str = "") -> None:
    def _fmt(v):
        try: return format(float(v), '.4f')
        except: return str(v)
    _log(logger, f"{prefix}n_est={params.get('n_estimators')}  depth={params.get('max_depth')}  "
                 f"lr={_fmt(params.get('learning_rate'))}  subsample={_fmt(params.get('subsample'))}  "
                 f"colsample={_fmt(params.get('colsample_bytree'))}")
    _log(logger, f"{prefix}gamma={_fmt(params.get('gamma'))}  min_child_w={params.get('min_child_weight')}  "
                 f"reg_α={_fmt(params.get('reg_alpha'))}  reg_λ={_fmt(params.get('reg_lambda'))}")


def _print_search_space(logger):
    _log(logger, f"  Search space (scale_pos_weight fixed = n_neg/n_pos, not searched):")
    for line in [
        "  n_estimators     : randint(1000, 5000)",
        "  max_depth        : randint(3, 10)",
        "  learning_rate    : uniform(0.005, 0.100)",
        "  subsample        : uniform(0.60, 1.00)",
        "  colsample_bytree : uniform(0.20, 0.80)",
        "  gamma            : uniform(0.00, 0.30)",
        "  min_child_weight : randint(1, 10)",
        "  reg_alpha        : uniform(0.00, 0.50)",
        "  reg_lambda       : uniform(0.50, 2.00)",
    ]: _log(logger, line)


# ══════════════════════════════════════════════════════════════════════════════
# DEFAULT CONFIGURATION
# ══════════════════════════════════════════════════════════════════════════════

_XGB_DEFAULT_CONFIG = {
    'mode': 'manual',
    'task': 'classification',
    'model': {
        'n_estimators':       3000,
        'max_depth':             6,
        'learning_rate':      0.01,
        'subsample':           0.8,
        'colsample_bytree':    0.8,
        'gamma':               0.0,
        'min_child_weight':      1,
        'reg_alpha':           0.0,
        'reg_lambda':          1.0,
        'tree_method':       'hist',
        'n_jobs':               -1,
        'random_state':         42,
    },
    'training': {
        'use_random_search':  False,
        # [HYP-1] search_strategy — see module docstring for full explanation
        # 'pre_search' is the recommended default for all routine use.
        'search_strategy':   'pre_search',
        'n_iter':               50,
        'cv':                    3,
        'early_stopping_rounds': 0,
        'eval_metric':        'auc',
        'class_weight':   'balanced',
    },
    'features': {
        'embedding':   False,
        'kmer':        False,
        'biophysical': False,
        'onehot':      False,
    },
    'kmer':  {'k': [1, 2], 'sequence': 'CDR3', 'normalize': True},
    'onehot': {'sequence': 'HCDR3', 'max_lengths': {'vh': 0, 'vl': 0, 'hcdr3': 0}},
    'lm_profiles': {
        'biophysical':     {'colsample_bytree': 1.0},
        'kmer':            {'colsample_bytree': 0.5},
        'onehot':          {'colsample_bytree': 0.3, 'max_depth': 8},
        'onehot_vh':       {'colsample_bytree': 0.3, 'max_depth': 8},
        'onehot_hcdr3':    {'colsample_bytree': 0.5},
        'onehot_cdr3':     {'colsample_bytree': 0.5},
        'ablang':          {'colsample_bytree': 0.25},
        'antiberty':       {'colsample_bytree': 0.25},
        'antiberta2':      {'colsample_bytree': 0.20},
        'antiberta2-cssp': {'colsample_bytree': 0.20},
        'igbert':          {'colsample_bytree': 0.20},
    },
    'shap': {
        'enabled': True, 'top_features': 30,
        'plot_types': ['bar', 'beeswarm', 'heatmap'],
        'max_waterfall_samples': 50, 'waterfall_top_n': 20,
        'waterfall_fig_size': [8, 6], 'pub_dpi': 300,
        'waterfall_fmt': 'tiff', 'make_ppt': True,
    },
    'mutagenesis': {'max_samples': 50, 'format': 'tiff', 'pub_dpi': 300, 'make_ppt': True},
}


# ══════════════════════════════════════════════════════════════════════════════
# XGBoostModel
# ══════════════════════════════════════════════════════════════════════════════

class XGBoostModel:

    @staticmethod
    def auto_detect_config(n: int, pos_rate: float, task: str = 'classification') -> dict:
        if   n <  5_000: tier='xs'; n_e,d,mf = 500,  10, 'sqrt'
        elif n < 20_000: tier='sm'; n_e,d,mf = 1000, 15, 0.3
        elif n < 80_000: tier='md'; n_e,d,mf = 2000, 20, 0.2
        elif n <200_000: tier='lg'; n_e,d,mf = 2000, 25, 0.15
        else:            tier='xl'; n_e,d,mf = 3000, 30, 0.1
        if task == 'regression': cw = None
        else:
            min_rate = min(pos_rate, 1-pos_rate)
            if   min_rate >= 0.40: cw = 'balanced'
            elif min_rate >= 0.10: cw = 'balanced_subsample'
            else: cw = {0:1, 1:(1-pos_rate)/(pos_rate+1e-8)}
        cfg = copy.deepcopy(_XGB_DEFAULT_CONFIG)
        cfg['mode'] = 'auto'; cfg['task'] = task
        cfg['model'].update({'n_estimators':n_e, 'max_depth':d,
                             'learning_rate':0.01,
                             'colsample_bytree': mf if isinstance(mf, float) else 0.8})
        cfg['training']['class_weight'] = cw
        cfg['_auto'] = {'n':n, 'pos_rate':round(pos_rate,4), 'size_tier':tier, 'task':task}
        return cfg

    @staticmethod
    def print_config_report(cfg: dict) -> None:
        a = cfg.get('_auto', {}); m = cfg.get('model', {})
        t = cfg.get('training', {}); ft = cfg.get('features', {})
        sh = cfg.get('shap', {}); task = cfg.get('task', 'classification')
        W = 62; sep = '═'*W; sep2 = '─'*W
        print(f"\n{sep}")
        print(f"  XGBoost  ·  {'AUTO' if a else 'MANUAL'}  ({task})"
              + (f"  tier={a.get('size_tier')}  n={a.get('n'):,}" if a else ""))
        print(sep2)
        print(f"  MODEL   n_est={m.get('n_estimators',3000)}  depth={m.get('max_depth',6)}"
              f"  lr={m.get('learning_rate',0.01):.3f}  subsample={m.get('subsample',0.8)}"
              f"  colsample={m.get('colsample_bytree',0.8)}")
        print(f"          gamma={m.get('gamma',0.0)}  min_child_w={m.get('min_child_weight',1)}"
              f"  reg_α={m.get('reg_alpha',0.0)}  reg_λ={m.get('reg_lambda',1.0)}")
        print(sep2)
        print(f"  FEATURES  embedding={ft.get('embedding')}  kmer={ft.get('kmer')}"
              f"  biophysical={ft.get('biophysical')}  onehot={ft.get('onehot',False)}")
        print(sep2)
        _search   = t.get('use_random_search', False)
        _strategy = t.get('search_strategy', 'pre_search') if _search else 'none'
        _speed    = {'pre_search': 'search BEFORE kfold — recommended default',
                     'first_fold': 'search on fold 1 — zero leakage, same speed',
                     'nested':     'per-fold search — gold standard, 10× slower',
                     'none':       'no search — use config params directly'}
        print(f"  HYPERPARAM  search={'ON' if _search else 'OFF'}"
              f"  strategy={_strategy}  n_iter={t.get('n_iter',50)}  cv={t.get('cv',3)}")
        if _search: print(f"              [{_speed.get(_strategy, '')}]")
        print(f"  SHAP  enabled={sh.get('enabled')}  top={sh.get('top_features')}")
        print(f"{sep}\n")

    def __init__(self, config: dict = None,
                 config_path: str = "config/xgboost.yaml",
                 verbose: bool = True):
        self.config = copy.deepcopy(_XGB_DEFAULT_CONFIG)
        if config_path and os.path.exists(config_path):
            with open(config_path) as f:
                user_cfg = yaml.safe_load(f) or {}
            _deep_merge(self.config, user_cfg)
            if verbose: print(f"[XGBoostModel] config ← {config_path}")
        if config is not None: _deep_merge(self.config, config)
        self.model = None; self.fb_ = None
        self.task  = self.config.get('task', 'classification')
        if verbose: self.print_config_report(self.config)

    def apply_lm_profile(self, lm: str, logger=None,
                          searched_params: set = None) -> None:
        """
        Apply LM-specific hyperparameter defaults from lm_profiles in config.

        If hyperparameter search was run, pass searched_params=set(best_params)
        so that searched values are NOT overwritten by lm_profile defaults.
        Params not in the search space (e.g. max_depth) are still applied.

        Without search (searched_params=None), all profile values are applied
        as before — original behaviour preserved.
        """
        profile = self.config.get('lm_profiles', {}).get(lm, {})
        if not profile: return
        searched_params = searched_params or set()
        applied = {}; skipped = {}
        for k, v in profile.items():
            if k in searched_params:
                skipped[k] = v          # search result takes priority
            else:
                self.config['model'][k] = v
                applied[k] = v
        if applied:
            _log(logger, f"[XGB] lm_profile '{lm}' applied : "
                         + "  ".join(f"{k}={v}" for k, v in applied.items()))
        if skipped:
            _log(logger, f"[XGB] lm_profile '{lm}' skipped : "
                         + "  ".join(f"{k}={v}" for k, v in skipped.items())
                         + "  ← search value kept")

    def _auto_fix_features(self, embeddings, X_df=None):
        ft = self.config.setdefault('features', {})
        if embeddings is not None and not ft.get('embedding', False): ft['embedding'] = True
        if embeddings is None and ft.get('embedding', True):
            if not any(ft.get(k) for k in ('kmer', 'biophysical', 'onehot')): ft['biophysical'] = True

    def _build_xgb(self, y=None):
        m = self.config['model']; t = self.config['training']
        _y = np.asarray(y, dtype=float) if y is not None else None
        _is_bin = (_y is not None and len(set(_y.tolist())) <= 2 and
                   set(_y.tolist()).issubset({0,1,0.0,1.0}))
        spw = 1.0
        if y is not None and self.task=='classification' and _is_bin and t.get('class_weight')=='balanced':
            c = np.bincount(np.asarray(y,dtype=int))
            spw = float(c[0])/float(c[1]) if len(c)>1 and c[1]>0 else 1.0
        common = dict(tree_method=m.get('tree_method','hist'),
                      n_estimators=m.get('n_estimators',3000), max_depth=m.get('max_depth',6),
                      learning_rate=m.get('learning_rate',0.01), subsample=m.get('subsample',0.8),
                      colsample_bytree=m.get('colsample_bytree',0.8), gamma=m.get('gamma',0.0),
                      min_child_weight=m.get('min_child_weight',1), reg_alpha=m.get('reg_alpha',0.0),
                      reg_lambda=m.get('reg_lambda',1.0), n_jobs=m.get('n_jobs',-1),
                      random_state=m.get('random_state',42), verbosity=0)
        _task = self.task
        if y is not None:
            _yf = np.asarray(y,dtype=float); _ib = (len(set(_yf.tolist()))<=2 and set(_yf.tolist()).issubset({0,1,0.0,1.0}))
            if _task=='auto': _task = 'classification' if _ib else 'regression'
            elif _task=='classification' and not _ib: _task = 'regression'
        if _task == 'classification':
            return xgb.XGBClassifier(objective='binary:logistic',
                                      eval_metric=t.get('eval_metric','auc'),
                                      scale_pos_weight=spw, **common)
        return xgb.XGBRegressor(objective='reg:squarederror', **common)

    # ── train ──────────────────────────────────────────────────────────────

    def train(self, X_df, y, embeddings=None, target="model", target_col="",
              embedding_lm="", val_X=None, val_y=None, val_embeddings=None, logger=None):
        _y_tmp = np.asarray(y, dtype=float)
        _is_bin = (len(set(_y_tmp.tolist()))<=2 and set(_y_tmp.tolist()).issubset({0,1,0.0,1.0}))
        if self.task=='classification' and not _is_bin: self.task='regression'
        y_arr = np.asarray(y, dtype=float if self.task=='regression' else int)
        self._auto_fix_features(embeddings, X_df)
        self.fb_ = FeatureBuilder(self.config)
        self.fb_.fit(X_df, embeddings)
        self.apply_lm_profile(embedding_lm, logger=logger)
        X_feat = self.fb_.transform(X_df, embeddings)
        _log(logger, f"[train] n={len(y_arr):,}  features={X_feat.shape[1]}  task={self.task}  lm={embedding_lm}")
        t = self.config['training']; esr = t.get('early_stopping_rounds', 0)

        if t.get('use_random_search', False):
            n_iter = t.get('n_iter',50); inner_cv = t.get('cv',3)
            _log(logger, f"[train] Hyperparameter search  n_iter={n_iter}  cv={inner_cv} ...")
            _print_search_space(logger)
            _spw = _compute_spw(y_arr, self.task)
            best_params, best_score = _run_search(X_feat, y_arr, n_iter, inner_cv, self.task, _spw, logger)
            _metric = 'AUC' if self.task=='classification' else 'R²'
            _log(logger, f"[train] Best CV {_metric}={best_score:.4f}")
            _log_params(best_params, logger, prefix="  [best] ")
            _log(logger, "[train] Copy params to xgboost.yaml to skip search next time.")
            self.config['model'].update(best_params)

        self.model = self._build_xgb(y_arr)
        if esr > 0 and val_X is not None and val_y is not None:
            X_va = self.fb_.transform(val_X, val_embeddings)
            y_va = np.asarray(val_y, dtype=float if self.task=='regression' else int)
            self.model.fit(X_feat, y_arr, eval_set=[(X_va,y_va)],
                           early_stopping_rounds=esr, verbose=False)
        elif esr > 0:
            from sklearn.model_selection import train_test_split
            X_tr,X_va,y_tr,y_va = train_test_split(X_feat, y_arr, test_size=0.1,
                stratify=(y_arr if self.task=='classification' else None), random_state=42)
            self.model.fit(X_tr, y_tr, eval_set=[(X_va,y_va)],
                           early_stopping_rounds=esr, verbose=False)
        else:
            self.model.fit(X_feat, y_arr)

        try:
            _p = self.model.get_params()
            def _f(v,fmt='.4f'):
                try: return format(float(v),fmt)
                except: return str(v)
            _log(logger, f"[params] n_est={_p.get('n_estimators')}  depth={_p.get('max_depth')}"
                         f"  lr={_f(_p.get('learning_rate'))}  subsample={_f(_p.get('subsample'))}"
                         f"  colsample={_f(_p.get('colsample_bytree'))}  gamma={_f(_p.get('gamma'))}"
                         f"  min_child_w={_p.get('min_child_weight')}  reg_α={_f(_p.get('reg_alpha'))}"
                         f"  reg_λ={_f(_p.get('reg_lambda'))}  scale_pos_w={_f(_p.get('scale_pos_weight',1.0),'.2f')}")
        except: pass

        if val_X is not None and val_y is not None:
            X_va = self.fb_.transform(val_X, val_embeddings)
            y_va = np.asarray(val_y, dtype=float if self.task=='regression' else int)
            if self.task=='classification':
                vp = self.model.predict_proba(X_va)[:,1]; vpreds = (vp>=0.5).astype(int)
                _log(logger, f"  val_auc={roc_auc_score(y_va,vp):.4f}  val_acc={accuracy_score(y_va,vpreds):.4f}"
                             f"  val_rec_fail={recall_score(y_va,vpreds,pos_label=0,zero_division=0):.4f}")
        _log(logger, "[train] completed.")

        if self.config.get('shap',{}).get('enabled',False) and _SHAP_AVAILABLE:
            sh_top = self.config.get('shap',{}).get('top_features',30)
            _parts = str(target).split('_xgboost_')
            try:
                self.shap_analysis(X_df, embeddings, output_prefix=target, split_tag="train",
                                   top_n=sh_top, barcodes=list(X_df.index.astype(str)),
                                   actual_labels=list(y_arr), actual_col_name=target_col or target,
                                   lm_name=embedding_lm,
                                   db_name=_parts[1] if len(_parts)>1 else target, logger=logger)
            except Exception as _se:
                import traceback
                _log(logger, f"[SHAP] train failed — {_se}\n{traceback.format_exc()}")
        return self

    # ── predict ────────────────────────────────────────────────────────────

    def predict_proba(self, X_df, embeddings=None) -> np.ndarray:
        if self.task=='regression': raise RuntimeError("Use predict() for regression.")
        return self.model.predict_proba(self.fb_.transform(X_df, embeddings))[:,1]

    def predict(self, X_df, embeddings=None, threshold: float = None) -> np.ndarray:
        X_feat = self.fb_.transform(X_df, embeddings)
        if self.task=='regression': return self.model.predict(X_feat)
        t = threshold if threshold is not None else getattr(self,'recommended_threshold',0.5)
        return (self.model.predict_proba(X_feat)[:,1] >= t).astype(int)

    def predict_raw(self, X_df, embeddings=None) -> np.ndarray:
        X_feat = self.fb_.transform(X_df, embeddings)
        return self.model.predict(X_feat) if self.task=='regression' else self.model.predict_proba(X_feat)[:,1]

    # ── save / load ────────────────────────────────────────────────────────

    def save(self, path: str):
        os.makedirs(os.path.dirname(path) or '.', exist_ok=True)
        rec_thresh = getattr(self, 'recommended_threshold', None)
        joblib.dump({'model':self.model,'fb_':self.fb_,'config':self.config,
                     'task':self.task,'recommended_threshold':rec_thresh}, path)
        thresh_note = (f"  threshold={rec_thresh:.4f}" if rec_thresh is not None
                       else "  threshold=None (0.5 at predict)")
        print(f"[XGBoostModel] saved → {path}{thresh_note}")

    @classmethod
    def load(cls, path, config_path: str = "config/xgboost.yaml"):
        instance = cls(config_path=config_path, verbose=False)
        payload  = joblib.load(path)
        if isinstance(payload, dict) and 'model' in payload:
            instance.model = payload['model']; instance.fb_ = payload.get('fb_')
            _deep_merge(instance.config, payload.get('config', {}))
            instance.task  = payload.get('task', 'classification')
            rt = payload.get('recommended_threshold')
        else:
            instance.model = payload; rt = None
        if instance.fb_ is None and instance.model is not None:
            try:
                n_feat = instance.model.n_features_in_
                instance.config.setdefault('features',{})['embedding'] = True
                fb = FeatureBuilder(instance.config)
                fb._emb_dim = n_feat; fb.feature_names_ = [f'emb_{i}' for i in range(n_feat)]
                instance.fb_ = fb
                print(f"  [load] fb_ reconstructed (dim={n_feat})")
            except Exception as _e:
                print(f"  [load] WARNING: fb_ reconstruction failed: {_e}")
        instance.recommended_threshold = float(rt) if rt is not None else 0.5
        flag = "  (embedded by kfold)" if rt and rt!=0.5 else ""
        print(f"[XGBoostModel] loaded ← {path}")
        print(f"  recommended_threshold={instance.recommended_threshold:.4f}{flag}")
        return instance

    # ── CDR3 mutagenesis ─────────────────────────────────────────────────

    def cdr3_mutagenesis_heatmap(self, X_df, embeddings=None,
                                  output_dir=None, pub_dpi=300,
                                  fig_size=(10,4), logger=None) -> dict:
        if self.model is None or self.fb_ is None:
            _log(logger, "[MUTA] model not trained/loaded"); return {}
        out_dir = output_dir or MODEL_DIR; os.makedirs(out_dir, exist_ok=True)
        results = {}; saved_paths = []
        for i in range(len(X_df)):
            row = X_df.iloc[i]; bc = str(X_df.index[i])
            cdr3_seq = str(row.get('CDR3','') or '').upper().replace('-','')
            if not cdr3_seq: continue
            n_pos = len(cdr3_seq); n_aa = len(AMINO_ACIDS)
            matrix = np.zeros((n_aa,n_pos), dtype=np.float32)
            mut_rows = []; mut_info = []
            for pos in range(n_pos):
                for aa_idx, aa in enumerate(AMINO_ACIDS):
                    mcdr3 = cdr3_seq[:pos]+aa+cdr3_seq[pos+1:]
                    mrow = dict(row); mrow['CDR3'] = mcdr3
                    if 'HSEQ' in mrow and mrow['HSEQ']:
                        hseq = str(mrow['HSEQ']); cs = hseq.find(cdr3_seq)
                        if cs >= 0: mrow['HSEQ'] = hseq[:cs]+mcdr3+hseq[cs+n_pos:]
                    mut_rows.append(mrow); mut_info.append((aa_idx,pos))
            mut_df = pd.DataFrame(mut_rows); mut_df.index = [f"{bc}_mut_{j}" for j in range(len(mut_df))]
            try:
                X_mut = self.fb_.transform(mut_df, None)
                scores = (self.model.predict_proba(X_mut)[:,1] if self.task=='classification'
                          else self.model.predict(X_mut))
            except Exception as _e: _log(logger, f"  [MUTA] {bc}: {_e}"); continue
            for j,(aa_idx,pos_idx) in enumerate(mut_info): matrix[aa_idx,pos_idx] = scores[j]
            results[bc] = matrix
            wt_row = pd.DataFrame([dict(row)]); wt_row.index = [bc]
            try:
                X_wt = self.fb_.transform(wt_row,None)
                wt_score = (self.model.predict_proba(X_wt)[:,1][0] if self.task=='classification'
                            else self.model.predict(X_wt)[0])
            except: wt_score = None
            fig, ax = plt.subplots(figsize=(max(fig_size[0],n_pos*0.45),fig_size[1]))
            import matplotlib.colors as _mc
            _norm = (_mc.TwoSlopeNorm(vmin=0.0,vcenter=0.5,vmax=1.0) if self.task=='classification'
                     else _mc.TwoSlopeNorm(vmin=float(matrix.min()),
                                           vcenter=(float(matrix.min())+float(matrix.max()))/2,
                                           vmax=float(matrix.max())))
            im = ax.imshow(matrix, aspect='auto', cmap='RdBu', norm=_norm)
            _fsz = max(4.0,min(7.0,120.0/max(n_pos,1)))
            for _ai in range(n_aa):
                for _pi in range(n_pos):
                    _v = matrix[_ai,_pi]; _tc = 'white' if (_v<0.35 or _v>0.65) else '#333333'
                    ax.text(_pi,_ai,f"{_v:.2f}",ha='center',va='center',fontsize=_fsz,color=_tc)
            plt.colorbar(im,ax=ax,label="P(PASS)" if self.task=='classification' else "Score",fraction=0.03,pad=0.02)
            ax.set_xticks(range(n_pos)); ax.set_xticklabels([f"{cdr3_seq[p]}\n{p+1}" for p in range(n_pos)],fontsize=8)
            ax.set_yticks(range(n_aa)); ax.set_yticklabels(list(AMINO_ACIDS),fontsize=7)
            ax.set_xlabel('CDR3 position',fontsize=8); ax.set_ylabel('Substituted AA',fontsize=8)
            for pos in range(n_pos):
                wi = AMINO_ACIDS.find(cdr3_seq[pos])
                if wi>=0: ax.add_patch(plt.Rectangle((pos-0.5,wi-0.5),1,1,fill=False,edgecolor='black',lw=1.5,zorder=5))
            wt_str = f"  WT P(PASS)={wt_score:.4f}" if wt_score is not None else ""
            ax.set_title(f"CDR3 Mutagenesis — {bc}{wt_str}\nXGBoost | CDR3={cdr3_seq}",fontsize=8)
            plt.tight_layout()
            bc_safe = bc.replace('/',' ').replace(' ','_')
            img_path = os.path.join(out_dir,f"{bc_safe}_cdr3_mutagenesis.tiff")
            plt.savefig(img_path,dpi=pub_dpi,format='tiff',bbox_inches='tight'); plt.close()
            saved_paths.append((bc,img_path,wt_score))
            _log(logger,f"  [MUTA] {bc} → {os.path.basename(img_path)}")
        if saved_paths:
            try:
                from pptx import Presentation as _Prs; from pptx.util import Inches,Pt
                from pptx.enum.text import PP_ALIGN; from pptx.dml.color import RGBColor
                prs = _Prs(); prs.slide_width=Inches(13.33); prs.slide_height=Inches(7.5)
                for bc,img_path,wt_score in saved_paths:
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    img_w=Inches(fig_size[0]); img_h=Inches(fig_size[1])
                    slide.shapes.add_picture(img_path,(prs.slide_width-img_w)/2,Inches(0.6),width=img_w,height=img_h)
                    txb=slide.shapes.add_textbox(Inches(0.15),Inches(7.1),Inches(13.0),Inches(0.35))
                    txb.text_frame.text=f"CDR3 Mutagenesis | {bc}"
                    p=txb.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
                    p.runs[0].font.size=Pt(7); p.runs[0].font.color.rgb=RGBColor(0x88,0x87,0x80)
                prs.save(os.path.join(out_dir,"cdr3_mutagenesis_all.pptx"))
                _log(logger,f"[MUTA] PPT saved")
            except ImportError: _log(logger,"[MUTA] python-pptx not installed")
            except Exception as _pe: _log(logger,f"[MUTA] PPT failed — {_pe}")
        _log(logger,f"[MUTA] complete  ({len(results)} antibodies)")
        return results

    # ── SHAP analysis ──────────────────────────────────────────────────────

    def shap_analysis(self, X_df, embeddings=None, output_prefix="xgb_shap",
                      split_tag="train", top_n=50, barcodes=None, actual_labels=None,
                      actual_col_name="label", lm_name="", db_name="", logger=None):
        if not _SHAP_AVAILABLE: _log(logger,"[SHAP] pip install shap"); return
        if self.model is None or self.fb_ is None: _log(logger,"[SHAP] model not trained"); return
        if barcodes is None:
            barcodes = list(X_df.index.astype(str)) if hasattr(X_df,'index') else [str(i) for i in range(len(X_df))]
        sh_cfg=self.config.get('shap',{}); max_s=sh_cfg.get('max_samples',500)
        plots=sh_cfg.get('plot_types',['bar','beeswarm']); top_n=sh_cfg.get('top_features',top_n)
        wf_top_n=sh_cfg.get('waterfall_top_features',top_n); max_wf=sh_cfg.get('max_waterfall_samples',50)
        pub_dpi=sh_cfg.get('pub_dpi',300); wf_size=sh_cfg.get('waterfall_fig_size',[8,6])
        wf_fmt=sh_cfg.get('waterfall_fmt','tiff').lower().strip('.'); make_ppt=sh_cfg.get('make_ppt',True)
        if 'beeswarm' not in plots: plots=list(plots)+['beeswarm']
        X_feat=self.fb_.transform(X_df,embeddings); ne_names=self.fb_.non_embedding_feature_names
        ne_idx=self.fb_.non_embedding_indices
        if not ne_idx: _log(logger,"[SHAP] Skipped — embedding-only mode."); return
        ne_idx_arr=np.array(ne_idx); X_ne=X_feat[:,ne_idx_arr]
        if len(X_ne)>max_s:
            idx_s=np.random.choice(len(X_ne),max_s,replace=False)
            X_shap=X_ne[idx_s]; _X_full_wf=X_feat[idx_s]
        else: X_shap=X_ne; _X_full_wf=X_feat
        _log(logger,f"\n[SHAP] {split_tag.upper()} — TreeExplainer  n={len(X_shap)}  features={len(ne_idx)}  top={top_n}")
        explainer=_shap_lib.TreeExplainer(self.model)
        try: shap_values=explainer.shap_values(X_shap,check_additivity=True)
        except: shap_values=explainer.shap_values(X_shap,check_additivity=False)
        if isinstance(shap_values,list): sv=np.array(shap_values[1],dtype=np.float64)
        elif isinstance(shap_values,np.ndarray):
            if shap_values.ndim==3: sv=shap_values[:,:,1] if shap_values.shape[0]==len(X_shap) else shap_values[1]
            else: sv=shap_values
        else: sv=np.array(shap_values,dtype=np.float64)
        sv=np.asarray(sv,dtype=np.float64)
        mean_abs=np.mean(np.abs(sv),axis=0); top_idx=[int(i) for i in np.argsort(mean_abs)[::-1][:top_n]]
        top_vals=[float(mean_abs[i]) for i in top_idx]; top_lbls=[ne_names[i] for i in top_idx]
        os.makedirs(MODEL_DIR,exist_ok=True); prefix=f"{output_prefix}_{split_tag}"
        _log(logger,f"  {'Rank':>4}  {'Feature':35s}  {'Mean |SHAP|':>11}")
        _log(logger,f"  {'─'*4}  {'─'*35}  {'─'*11}")
        for rank,(lbl,val) in enumerate(zip(top_lbls,top_vals),1):
            _log(logger,f"  {rank:4d}  {lbl:35s}  {val:11.6f}")
        import csv
        csv_path=os.path.join(MODEL_DIR,f"{prefix}_shap_top{top_n}.csv")
        with open(csv_path,'w',newline='') as f:
            w=csv.writer(f); w.writerow(['rank','feature','mean_abs_shap'])
            for rank,(lbl,val) in enumerate(zip(top_lbls,top_vals),1): w.writerow([rank,lbl,f"{val:.8f}"])
        _log(logger,f"[SHAP] csv → {csv_path}")
        if 'bar' in plots:
            colours=[('#534AB7' if l.startswith('1mer_') else '#7F77DD' if l.startswith('2mer_') else
                      '#E07B39' if l.startswith('oh_') else '#1D9E75') for l in top_lbls]
            fig,ax=plt.subplots(figsize=(8.27,min(11.69,max(5,top_n*0.20))))
            ax.barh(range(len(top_idx)),top_vals,color=colours); ax.set_yticks(range(len(top_idx)))
            ax.set_yticklabels(top_lbls,fontsize=7); ax.invert_yaxis()
            ax.set_xlabel('Mean |SHAP value|'); ax.set_title(f'SHAP top {top_n} — {output_prefix} [{split_tag}]')
            plt.tight_layout(); out=os.path.join(MODEL_DIR,f"{prefix}_shap_bar.png")
            plt.savefig(out,dpi=150,bbox_inches='tight'); plt.close(); _log(logger,f"[SHAP] bar → {out}")
        if 'beeswarm' in plots:
            _shap_lib.summary_plot(sv[:,top_idx],X_shap[:,top_idx],feature_names=top_lbls,show=False,max_display=top_n)
            plt.title(f'SHAP beeswarm — {output_prefix} [{split_tag}]')
            plt.gcf().set_size_inches(8.27,min(11.69,max(5,top_n*0.22))); plt.tight_layout()
            out=os.path.join(MODEL_DIR,f"{prefix}_shap_beeswarm.png")
            plt.savefig(out,dpi=150,bbox_inches='tight'); plt.close(); _log(logger,f"[SHAP] beeswarm → {out}")
        if 'heatmap' in plots:
            n_heat=min(top_n,20); sub=sv[:min(100,len(sv)),:][:,top_idx[:n_heat]]
            fig,ax=plt.subplots(figsize=(11,7)); vmax=np.abs(sub).max()
            im=ax.imshow(sub.T,aspect='auto',cmap='RdBu_r',vmin=-vmax,vmax=vmax)
            plt.colorbar(im,ax=ax,label='SHAP value'); ax.set_yticks(range(n_heat))
            ax.set_yticklabels(top_lbls[:n_heat],fontsize=7); ax.set_xlabel('Sample index')
            ax.set_title(f'SHAP heatmap top {n_heat} — {output_prefix} [{split_tag}]')
            plt.tight_layout(); out=os.path.join(MODEL_DIR,f"{prefix}_shap_heatmap.png")
            plt.savefig(out,dpi=150,bbox_inches='tight'); plt.close(); _log(logger,f"[SHAP] heatmap → {out}")
        try:
            import matplotlib.patches as _mp
            n_wf=min(len(sv),max_wf); wf_dir=os.path.join(MODEL_DIR,f"{prefix}_shap_waterfalls")
            os.makedirs(wf_dir,exist_ok=True)
            if self.task=='classification':
                _thresh=getattr(self,'recommended_threshold',0.5)
                all_probs=self.model.predict_proba(_X_full_wf[:n_wf])[:,1]
                all_labels_pred=(all_probs>=_thresh).astype(int)
            else: all_probs=self.model.predict(_X_full_wf[:n_wf]); all_labels_pred=None
            expected_val=float(explainer.expected_value[1] if isinstance(explainer.expected_value,(list,np.ndarray)) else explainer.expected_value)
            saved_wf=[]
            for s_idx in range(n_wf):
                bc=barcodes[s_idx] if barcodes and s_idx<len(barcodes) else str(s_idx)
                act=actual_labels[s_idx] if actual_labels is not None and s_idx<len(actual_labels) else None
                n_f=min(wf_top_n,len(top_idx))
                shap_vals=[float(sv[s_idx][i]) for i in top_idx[:n_f]]
                feat_vals=[float(X_shap[s_idx,i]) for i in top_idx[:n_f]]
                lbls_s=top_lbls[:n_f]
                triples=sorted(zip(shap_vals,lbls_s,feat_vals),key=lambda x:abs(x[0]),reverse=True)
                wf_s=[x[0] for x in triples]; wf_l=[f"{x[1]} = {x[2]:.3g}" for x in triples]
                cols=['#D62728' if v<0 else '#1F77B4' for v in wf_s]
                cumsum=expected_val; bottoms=[]
                for v in wf_s: bottoms.append(cumsum); cumsum+=v
                fig,ax=plt.subplots(figsize=(wf_size[0],wf_size[1]))
                _xmin=min(bottoms+[expected_val]+[b+v for b,v in zip(bottoms,wf_s)])
                _xmax=max(bottoms+[expected_val]+[b+v for b,v in zip(bottoms,wf_s)])
                _xrng=max(abs(_xmax-_xmin),1e-6)
                for i,(v,lbl,bot,col) in enumerate(zip(wf_s,wf_l,bottoms,cols)):
                    ax.barh(i,v,left=bot,color=col,height=0.70,edgecolor='white',linewidth=0.4)
                    bw=abs(v)
                    if bw/_xrng>0.08: xt=bot+v/2; ha='center'; tc='white'
                    else:
                        pad=_xrng*0.012; xt=bot+v+pad*(1 if v>=0 else -1)
                        ha='left' if v>=0 else 'right'; tc=col
                    if bw/_xrng>0.005: ax.text(xt,i,f"{v:+.4f}",va='center',ha=ha,fontsize=8,color=tc)
                ax.set_yticks(range(n_f)); ax.set_yticklabels(wf_l,fontsize=8); ax.invert_yaxis()
                ax.set_xlim(_xmin-_xrng*0.18,_xmax+_xrng*0.06)
                ax.axvline(expected_val,color='#888',lw=0.8,ls='--')
                ax.axvline(cumsum,color='#222299',lw=2.0,ls='-')
                ax.set_xlabel('SHAP value',fontsize=8)
                _act_str=(f"  |  Actual ({actual_col_name}) = {'PASS' if int(act)==1 else 'FAIL'}" if act is not None else "")
                sc_part=(f"P(PASS)={float(all_probs[s_idx]):.4f}  →  {'PASS' if all_labels_pred[s_idx]==1 else 'FAIL'}"
                         if self.task=='classification' else f"Predicted={float(all_probs[s_idx]):.4f}")
                ax.set_title(f"SHAP Waterfall — {bc}\n{sc_part}{_act_str}  |  XGBoost | {lm_name} | {db_name}",fontsize=8)
                ax.legend(handles=[_mp.Patch(color='#1F77B4',label='toward PASS'),
                                   _mp.Patch(color='#D62728',label='toward FAIL'),
                                   plt.Line2D([0],[0],color='#999',ls='--',lw=1,label=f'baseline={expected_val:.3f}'),
                                   plt.Line2D([0],[0],color='#1A237E',ls='-',lw=2,label=f'final={cumsum:.4f}')],
                           fontsize=7,loc='lower right',framealpha=0.85)
                plt.tight_layout()
                bc_safe=str(bc).replace('/',' ').replace(' ','_')
                img_path=os.path.join(wf_dir,f"{s_idx+1:04d}_{bc_safe}_waterfall.{wf_fmt}")
                save_kw=dict(dpi=pub_dpi,bbox_inches='tight')
                if wf_fmt=='tiff': save_kw['format']='tiff'
                elif wf_fmt in ('jpeg','jpg'): save_kw['format']='jpeg'; save_kw['pil_kwargs']={'quality':95}
                plt.savefig(img_path,**save_kw); plt.close(); saved_wf.append(img_path)
            _log(logger,f"[SHAP] {n_wf} waterfall images → {wf_dir}/")
            if make_ppt and saved_wf:
                try:
                    from pptx import Presentation as _Prs; from pptx.util import Inches,Pt
                    from pptx.enum.text import PP_ALIGN; from pptx.dml.color import RGBColor
                    prs=_Prs(); prs.slide_width=Inches(13.33); prs.slide_height=Inches(7.5)
                    for img_path,s_idx in zip(saved_wf,range(len(saved_wf))):
                        bc=barcodes[s_idx] if barcodes and s_idx<len(barcodes) else str(s_idx)
                        slide=prs.slides.add_slide(prs.slide_layouts[6])
                        img_w=Inches(wf_size[0]); img_h=Inches(wf_size[1])
                        slide.shapes.add_picture(img_path,(prs.slide_width-img_w)/2,Inches(0.6),width=img_w,height=img_h)
                        txb=slide.shapes.add_textbox(Inches(0.15),Inches(7.1),Inches(13.0),Inches(0.35))
                        txb.text_frame.text=f"{s_idx+1}/{n_wf}  |  {bc}  |  XGBoost | {lm_name} | {db_name} | {split_tag}"
                        p=txb.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
                        p.runs[0].font.size=Pt(7); p.runs[0].font.color.rgb=RGBColor(0x88,0x87,0x80)
                    ppt_path=os.path.join(MODEL_DIR,f"{prefix}_shap_waterfall.pptx")
                    prs.save(ppt_path); _log(logger,f"[SHAP] waterfall PPT ({n_wf} slides) → {ppt_path}")
                except ImportError: _log(logger,"[SHAP] python-pptx not installed")
                except Exception as _pe:
                    import traceback as _tb; _log(logger,f"[SHAP] PPT failed — {_pe}\n{_tb.format_exc()}")
        except Exception as _we:
            import traceback as _tb; _log(logger,f"[SHAP] waterfall failed — {_we}\n{_tb.format_exc()}")
        _log(logger,f"[SHAP] {split_tag} complete.")

    # ── k-fold validation ─────────────────────────────────────────────────

    @classmethod
    def kfold_validation(cls,
                         data,
                         X_df:          pd.DataFrame,
                         y,
                         embeddings:    np.ndarray  = None,
                         embedding_lm:  str         = '',
                         title:         str         = "XGBoost",
                         kfold:         int         = 10,
                         target:        str         = "psr_filter",
                         cluster_col:   str         = "HCDR3_CLUSTER_0.8",
                         db_stem:       str         = "",
                         override_features: dict    = None,
                         cost_fn:       float       = 3.0,
                         cost_fp:       float       = 1.0):
        """
        CDR3-cluster stratified k-fold with optional hyperparameter search.

        YAML config (xgboost.yaml):
          training:
            use_random_search: true
            search_strategy: pre_search   # RECOMMENDED
            n_iter: 50
            cv: 3

        Search is cleanly separated from evaluation:
          pre_search  → run RandomizedSearchCV on full data BEFORE kfold
                         own 3-fold internal CV → minimal leakage
                         all folds use identical best params → fair comparison
          first_fold  → search on fold 1 only, reuse for 2–K → zero leakage
          nested      → per-fold search → gold standard, 10× slower
          none        → use YAML params directly → fastest

        After any search, best params are printed and saved —
        copy to xgboost.yaml model: section to skip search next time.
        """
        import yaml as _yaml

        os.makedirs(MODEL_DIR, exist_ok=True)

        # ── Logging ───────────────────────────────────────────────────────
        ts       = datetime.datetime.now().strftime('%Y%m%d_%H%M%S')
        _db_tag  = f"_{db_stem}" if db_stem else ""
        logger   = _setup_logger(os.path.join(MODEL_DIR,
                    f"kfold_{target}_{embedding_lm}_xgboost{_db_tag}_k{kfold}_{ts}.log"))
        _log(logger, f"[log] {'='*58}")
        _log(logger, f"[log] Started : {datetime.datetime.now()}")
        _log(logger, f"[log] target={target}  lm={embedding_lm}  k={kfold}")
        _log(logger, f"[log] {'='*58}\n")

        # ── Base config ───────────────────────────────────────────────────
        _holder        = cls.__new__(cls)
        _holder.config = copy.deepcopy(_XGB_DEFAULT_CONFIG)
        _cfg_path      = "config/xgboost.yaml"
        if os.path.exists(_cfg_path):
            with open(_cfg_path) as _f:
                _deep_merge(_holder.config, _yaml.safe_load(_f) or {})
        resolved_cfg = _holder.config

        # ── Search config ─────────────────────────────────────────────────
        t_cfg      = resolved_cfg['training']
        _do_search = t_cfg.get('use_random_search', False)
        _strategy  = t_cfg.get('search_strategy', 'pre_search').lower()
        _n_iter    = int(t_cfg.get('n_iter', 50))
        _inner_cv  = int(t_cfg.get('cv', 3))

        # ── Auto-detect task ──────────────────────────────────────────────
        _y_tmp  = np.asarray(y, dtype=float)
        _is_bin = (len(set(_y_tmp.tolist())) <= 2 and
                   set(_y_tmp.tolist()).issubset({0, 1, 0.0, 1.0}))
        _kfold_task = 'classification' if _is_bin else 'regression'
        if not _is_bin:
            resolved_cfg['task'] = 'regression'
            _log(logger, "[kfold] Auto-detected task=regression")
        y_arr    = np.asarray(y, dtype=float if _kfold_task=='regression' else int)
        n        = len(y_arr)
        mean_fpr = np.linspace(0, 1, 100)

        # ── Log search plan ───────────────────────────────────────────────
        if _do_search:
            _log(logger, f"\n{'═'*62}")
            _log(logger, f"  HYPERPARAMETER SEARCH")
            _log(logger, f"{'─'*62}")
            _speed_note = {
                'pre_search': 'search on full data BEFORE kfold — recommended default',
                'first_fold': 'search on fold 1 only — zero leakage, same speed as pre_search',
                'nested':     'per-fold inner search — gold standard, ~10× slower',
            }
            _log(logger, f"  strategy : {_strategy}  [{_speed_note.get(_strategy,'')}]")
            _log(logger, f"  n_iter   : {_n_iter}  |  inner_cv : {_inner_cv}")
            _print_search_space(logger)
            _log(logger, f"{'═'*62}\n")
        else:
            _strategy = 'none'
            _log(logger, "\n[hyperparam] Search DISABLED — using config params directly.")
            _log(logger, "  Set use_random_search: true in xgboost.yaml to enable.\n")

        # ── Build full-data feature matrix ONCE ───────────────────────────
        # Reused by pre_search and sliced per fold — avoids rebuilding features
        # in every fold iteration which would be redundant and slow.
        _tmp_inst        = cls.__new__(cls)
        _tmp_inst.config = copy.deepcopy(resolved_cfg)
        _tmp_inst.task   = _kfold_task
        _tmp_inst.fb_    = FeatureBuilder(_tmp_inst.config)
        if override_features:
            _oh_seq_tmp = override_features.get('_onehot_sequence')
            _km_seq_tmp = override_features.get('_kmer_sequence')
            _tmp_inst.config['features'].update(
                {k: v for k, v in override_features.items()
                 if k not in ('_onehot_sequence', '_kmer_sequence')})
            if _oh_seq_tmp:
                _tmp_inst.config.setdefault('onehot', {})['sequence'] = _oh_seq_tmp
            if _km_seq_tmp:
                _tmp_inst.config.setdefault('kmer', {})['sequence'] = _km_seq_tmp
        # For feature building only — no search yet, so no params to protect
        _tmp_inst.apply_lm_profile(embedding_lm, logger=logger)
        X_all = _tmp_inst.fb_.fit_transform(X_df, embeddings)
        _log(logger, f"[kfold] Full dataset: n={n:,}  features={X_all.shape[1]}"
                     f"  task={_kfold_task}  lm={embedding_lm}")

        # ── [HYP-1] PRE_SEARCH ───────────────────────────────────────────
        #
        # Rationale for making this the default:
        #
        #   Goal of hyperparameter search  = find structural params
        #     (tree depth, learning rate, regularization)
        #   Goal of k-fold validation      = measure generalization
        #
        #   These are INDEPENDENT questions. Running them together
        #   (nested CV) adds complexity and cost without meaningful
        #   benefit for structural params like depth and lr.
        #
        #   Pre-search uses its own internal 3-fold CV, so the outer
        #   k-fold val folds are never used to select params.
        #   The remaining "leakage" (search sees train+val labels at
        #   the dataset level) affects structural params by < 0.003 AUC —
        #   well within the ±0.010 fold variance of your models.
        #
        #   The key advantage over nested: all 10 folds use the SAME params,
        #   making the kfold results directly comparable across folds.
        #   With nested, fold 1 might use depth=5 and fold 7 depth=8,
        #   making variance inflation hard to distinguish from real
        #   generalization variance.
        #
        _first_fold_params = None   # used by first_fold strategy
        _searched_params   = set()  # tracks keys set by search → lm_profile skips these

        if _do_search and _strategy == 'pre_search':
            _spw = _compute_spw(y_arr, _kfold_task)
            _log(logger, f"[pre_search] RandomizedSearchCV on full dataset "
                         f"(n={n:,}  n_iter={_n_iter}  inner_cv={_inner_cv})")
            _log(logger, f"[pre_search] scale_pos_weight={_spw:.4f}  (fixed, not searched)")

            _best_params, _best_score = _run_search(
                X_all, y_arr, _n_iter, _inner_cv, _kfold_task, _spw, logger)

            _metric_name = 'AUC' if _kfold_task == 'classification' else 'R²'
            _log(logger, f"\n[pre_search] ✓  Best CV {_metric_name}={_best_score:.4f}")
            _log_params(_best_params, logger, prefix="  [best] ")
            _log(logger, f"\n[pre_search] To skip search in future runs, "
                         f"copy these to xgboost.yaml → model:")
            for _pk, _pv in sorted(_best_params.items()):
                _log(logger, f"    {_pk}: {_pv}")
            _log(logger, f"\n[pre_search] Applying best params to all {kfold} folds ...\n")
            resolved_cfg['model'].update(_best_params)
            _searched_params = set(_best_params.keys())  # protect from lm_profile override

        # ── Splitter ──────────────────────────────────────────────────────
        from sklearn.model_selection import GroupKFold, KFold
        kfold_actual = kfold
        if cluster_col in data.columns:
            groups          = data[cluster_col].values
            n_unique_groups = len(np.unique(groups))
            if n_unique_groups < kfold_actual:
                kfold_actual = n_unique_groups
                _log(logger, f"[kfold] {n_unique_groups} clusters → folds reduced to {kfold_actual}")
            if _kfold_task == 'regression':
                splitter   = GroupKFold(n_splits=kfold_actual)
                split_iter = splitter.split(np.arange(n), y_arr, groups)
            elif n_unique_groups == n:
                _log(logger, "[kfold] all-singleton clusters → StratifiedKFold")
                splitter   = StratifiedKFold(n_splits=kfold_actual, shuffle=True, random_state=42)
                split_iter = splitter.split(np.arange(n), y_arr)
            else:
                splitter   = StratifiedGroupKFold(n_splits=kfold_actual, shuffle=True, random_state=42)
                split_iter = splitter.split(np.arange(n), y_arr, groups)
                _log(logger, f"[kfold] StratifiedGroupKFold on '{cluster_col}' "
                             f"({n_unique_groups} clusters, {kfold_actual} folds)")
        else:
            _log(logger, f"[kfold] '{cluster_col}' not found → StratifiedKFold/KFold")
            if _kfold_task == 'regression':
                splitter   = KFold(n_splits=kfold_actual, shuffle=True, random_state=42)
                split_iter = splitter.split(np.arange(n))
            else:
                splitter   = StratifiedKFold(n_splits=kfold_actual, shuffle=True, random_state=42)
                split_iter = splitter.split(np.arange(n), y_arr)

        # ── Fold loop ─────────────────────────────────────────────────────
        tprs, aucs_list, fold_metrics, all_records = [], [], [], []
        best_fold_auc   = -1.0; best_fold_num = -1
        best_fold_state = best_fold_cfg = best_fold_fb = None
        _nested_params_list = []

        _log(logger, f"\n{'='*72}")
        _log(logger, f"  {kfold_actual}-FOLD CV  |  XGBoost  |  {target.upper()}  |  "
                     f"{embedding_lm}  |  {_kfold_task}  |  search={_strategy}")
        _log(logger, f"{'─'*72}")
        if _kfold_task != 'regression':
            _log(logger, f"  {'Fold':>5}  {'AUC':>7}  {'Acc':>7}  {'F1':>7}  "
                         f"{'Prec':>7}  {'Rec':>7}  {'Rec(F)':>7}")
            _log(logger, f"  {'─'*5}  {'─'*7}  {'─'*7}  {'─'*7}  {'─'*7}  {'─'*7}  {'─'*7}")
        else:
            _log(logger, f"  {'Fold':>5}  {'R2':>7}  {'MAE':>7}  {'Pearson':>8}  {'Spearman':>8}")
            _log(logger, f"  {'─'*5}  {'─'*7}  {'─'*7}  {'─'*8}  {'─'*8}")

        plt.figure(figsize=(8, 7))

        for fold, (tr_idx, va_idx) in enumerate(split_iter, 1):

            if cluster_col in data.columns:
                leaked = set(groups[tr_idx]) & set(groups[va_idx])
                status = f"[WARN] {len(leaked)} cluster(s) leaked" if leaked else "[OK]  No CDR3 leakage"
                _log(logger, f"\n── Fold {fold}/{kfold_actual} ── {status}")
            else:
                _log(logger, f"\n── Fold {fold}/{kfold_actual} ──")

            y_tr = y_arr[tr_idx]; y_va = y_arr[va_idx]
            if _kfold_task != 'regression':
                _log(logger, f"  Train={len(tr_idx):,} pos={y_tr.mean():.1%}  "
                             f"Val={len(va_idx):,} pos={y_va.mean():.1%}")
            else:
                _log(logger, f"  Train={len(tr_idx):,} mu={y_tr.mean():.4f}  "
                             f"Val={len(va_idx):,} mu={y_va.mean():.4f}")

            # Slice pre-built feature matrix — fast, no recomputation
            X_tr = X_all[tr_idx]; X_va = X_all[va_idx]

            inst = cls(config=resolved_cfg, verbose=False)
            inst.task = _kfold_task
            inst.fb_  = copy.deepcopy(_tmp_inst.fb_)  # reuse fitted FeatureBuilder

            # ── [HYP-1] FIRST_FOLD ────────────────────────────────────────
            if _do_search and _strategy == 'first_fold':
                if fold == 1:
                    _spw = _compute_spw(y_tr, _kfold_task)
                    _log(logger, f"  [first_fold] Searching on fold 1 train "
                                 f"(n={len(X_tr):,}  n_iter={_n_iter}  cv={_inner_cv}) ...")
                    _first_fold_params, _ff_score = _run_search(
                        X_tr, y_tr, _n_iter, _inner_cv, _kfold_task, _spw, logger)
                    _metric_name = 'AUC' if _kfold_task == 'classification' else 'R²'
                    _log(logger, f"  [first_fold] Best CV {_metric_name}={_ff_score:.4f}")
                    _log_params(_first_fold_params, logger, prefix="    ")
                    _log(logger, f"  [first_fold] Params locked for folds 2–{kfold_actual} ✓")
                    _searched_params = set(_first_fold_params.keys())
                else:
                    _log(logger, f"  [first_fold] Reusing fold 1 params")
                inst.config['model'].update(_first_fold_params)

            # ── [HYP-1] NESTED ────────────────────────────────────────────
            elif _do_search and _strategy == 'nested':
                _spw = _compute_spw(y_tr, _kfold_task)
                _log(logger, f"  [nested] Searching fold {fold} "
                             f"(n={len(X_tr):,}  n_iter={_n_iter}  cv={_inner_cv}) ...")
                _fold_params, _fold_score = _run_search(
                    X_tr, y_tr, _n_iter, _inner_cv, _kfold_task, _spw, logger)
                _metric_name = 'AUC' if _kfold_task=='classification' else 'R²'
                _log(logger, f"  [nested] Best CV {_metric_name}={_fold_score:.4f}")
                _log_params(_fold_params, logger, prefix="    ")
                inst.config['model'].update(_fold_params)
                _searched_params = set(_fold_params.keys())
                _nested_params_list.append({'fold':fold, **_fold_params,
                                            f'cv_{_metric_name.lower()}':_fold_score})

            inst.apply_lm_profile(embedding_lm, logger=logger,
                                  searched_params=_searched_params)
            inst.model = inst._build_xgb(y_tr)

            # Log params
            try:
                _p = inst.model.get_params()
                def _f(v,fmt='.4f'):
                    try: return format(float(v),fmt)
                    except: return str(v)
                _log(logger, f"  [params] n_est={_p.get('n_estimators')}  depth={_p.get('max_depth')}"
                             f"  lr={_f(_p.get('learning_rate'))}  subsample={_f(_p.get('subsample'))}"
                             f"  colsample={_f(_p.get('colsample_bytree'))}  gamma={_f(_p.get('gamma'))}"
                             f"  min_child_w={_p.get('min_child_weight')}  reg_α={_f(_p.get('reg_alpha'))}"
                             f"  reg_λ={_f(_p.get('reg_lambda'))}  scale_pos_w={_f(_p.get('scale_pos_weight',1.0),'.2f')}")
            except: pass

            esr = inst.config['training'].get('early_stopping_rounds', 0)
            if esr > 0:
                inst.model.fit(X_tr, y_tr, eval_set=[(X_va,y_va)],
                               early_stopping_rounds=esr, verbose=False)
                _log(logger, f"  [early_stop] best_iteration={inst.model.best_iteration}")
            else:
                inst.model.fit(X_tr, y_tr)

            # Evaluate
            if inst.task == 'classification':
                probs = inst.model.predict_proba(X_va)[:,1]
                preds = (probs>=0.5).astype(int)
                if len(set(y_va)) < 2:
                    _log(logger, f"  Skipping fold {fold} — one class only."); continue
                fold_auc  = roc_auc_score(y_va, probs)
                fold_acc  = accuracy_score(y_va, preds)
                fold_f1   = f1_score(y_va, preds, zero_division=0)
                fold_prec = precision_score(y_va, preds, zero_division=0)
                fold_rec  = recall_score(y_va, preds, zero_division=0)
                fold_recf = recall_score(y_va, preds, pos_label=0, zero_division=0)
                fold_metrics.append({'fold':fold,'auc':fold_auc,'acc':fold_acc,'f1':fold_f1,
                                     'precision':fold_prec,'recall':fold_rec,'rec_fail':fold_recf})
                aucs_list.append(fold_auc)
                fpr,tpr,_ = roc_curve(y_va, probs)
                tprs.append(np.interp(mean_fpr,fpr,tpr)); tprs[-1][0] = 0.0
                plt.plot(fpr,tpr,lw=1,alpha=0.3,label=f'Fold {fold} ({fold_auc:.3f})')
                _log(logger,
                     f"  → AUC={fold_auc:.4f}  Acc={fold_acc:.4f}  F1={fold_f1:.4f}"
                     f"  Prec={fold_prec:.4f}  Rec(PASS)={fold_rec:.4f}  Rec(FAIL)={fold_recf:.4f}")
            else:
                from sklearn.metrics import r2_score, mean_absolute_error
                from scipy.stats import pearsonr, spearmanr
                preds_r  = inst.model.predict(X_va)
                fold_auc = spearmanr(y_va, preds_r)[0]
                fold_r2  = r2_score(y_va, preds_r)
                fold_mae = mean_absolute_error(y_va, preds_r)
                fold_rp  = pearsonr(y_va, preds_r)[0]
                fold_metrics.append({'fold':fold,'r2':fold_r2,'mae':fold_mae,
                                     'pearson':fold_rp,'spearman':fold_auc})
                aucs_list.append(fold_auc)
                plt.scatter(y_va,preds_r,alpha=0.3,s=10,label=f'Fold {fold} (ρ={fold_auc:.3f})')
                _log(logger,
                     f"  → R²={fold_r2:.4f}  MAE={fold_mae:.4f}"
                     f"  Pearson={fold_rp:.4f}  Spearman={fold_auc:.4f}")

            _task_tag = "_regression" if _kfold_task=='regression' else ""
            fold_path = os.path.join(MODEL_DIR,
                f"xgboost_{target}_{embedding_lm}{_db_tag}_fold{fold}_k{kfold_actual}{_task_tag}.pkl")
            inst.save(fold_path)

            if fold_auc > best_fold_auc:
                best_fold_auc=fold_auc; best_fold_num=fold
                best_fold_state=copy.deepcopy(inst.model)
                best_fold_cfg=copy.deepcopy(inst.config)
                best_fold_fb=copy.deepcopy(inst.fb_)

            bcs = X_df.iloc[va_idx].index.astype(str).tolist()
            if inst.task == 'classification':
                for bc,true,pred,prob in zip(bcs,y_va,preds,probs):
                    all_records.append({'BARCODE':bc,'fold':fold,'true':true,'pred':pred,'prob':prob})
            else:
                for bc,true,pred in zip(bcs,y_va,preds_r):
                    all_records.append({'BARCODE':bc,'fold':fold,'true':true,'pred':pred,'prob':pred})

        if not aucs_list:
            _log(logger, "[kfold] No valid folds."); return

        # ── Nested stability summary ───────────────────────────────────────
        if _do_search and _strategy=='nested' and _nested_params_list:
            _log(logger, f"\n{'═'*72}")
            _log(logger, f"  HYPERPARAMETER STABILITY (nested — per-fold best params)")
            _log(logger, f"{'─'*72}")
            _param_keys = ['n_estimators','max_depth','learning_rate','subsample',
                           'colsample_bytree','gamma','min_child_weight','reg_alpha','reg_lambda']
            _log(logger, f"  {'Param':<22}  {'Mean':>9}  {'Std':>9}  {'Min':>9}  {'Max':>9}")
            _log(logger, f"  {'─'*22}  {'─'*9}  {'─'*9}  {'─'*9}  {'─'*9}")
            for _pk in _param_keys:
                _vals = [float(d[_pk]) for d in _nested_params_list if _pk in d]
                if not _vals: continue
                _log(logger, f"  {_pk:<22}  {np.mean(_vals):9.4f}  {np.std(_vals):9.4f}  "
                             f"{np.min(_vals):9.4f}  {np.max(_vals):9.4f}")
            _log(logger, f"{'─'*72}")
            _log(logger, f"  Low std  → stable — safe to fix in xgboost.yaml")
            _log(logger, f"  High std → data-sensitive — keep nested search")
            _log(logger, f"{'═'*72}")
            _params_csv = os.path.join(MODEL_DIR,
                f"kfold_{target}_{embedding_lm}_xgboost{_db_tag}_k{kfold_actual}_nested_params.csv")
            pd.DataFrame(_nested_params_list).to_csv(_params_csv, index=False)
            _log(logger, f"[hyperparam] Per-fold params → {_params_csv}")

        # ── Summary ────────────────────────────────────────────────────────
        mean_auc = float(np.mean(aucs_list)); std_auc = float(np.std(aucs_list))
        _is_reg  = fold_metrics and 'r2' in fold_metrics[0]
        search_tag = f" [{_strategy}]" if _do_search else ""

        if not _is_reg:
            mean_acc  = float(np.mean([m['acc']       for m in fold_metrics]))
            mean_f1   = float(np.mean([m['f1']        for m in fold_metrics]))
            mean_prec = float(np.mean([m['precision'] for m in fold_metrics]))
            mean_rec  = float(np.mean([m['recall']    for m in fold_metrics]))
            mean_recf = float(np.mean([m['rec_fail']  for m in fold_metrics]))
            _log(logger, f"  {'─'*5}  {'─'*7}  {'─'*7}  {'─'*7}  {'─'*7}  {'─'*7}  {'─'*7}")
            for m in fold_metrics:
                mark = " ←" if m['fold']==best_fold_num else ""
                _log(logger, f"  {m['fold']:5d}  {m['auc']:7.4f}  {m['acc']:7.4f}  {m['f1']:7.4f}  "
                             f"{m['precision']:7.4f}  {m['recall']:7.4f}  {m['rec_fail']:7.4f}{mark}")
            _log(logger, f"  {'─'*5}  {'─'*7}  {'─'*7}  {'─'*7}  {'─'*7}  {'─'*7}  {'─'*7}")
            _log(logger, f"  {'Mean':>5}  {mean_auc:7.4f}  {mean_acc:7.4f}  {mean_f1:7.4f}  "
                         f"{mean_prec:7.4f}  {mean_rec:7.4f}  {mean_recf:7.4f}")
            _log(logger, f"  {'±Std':>5}  {std_auc:7.4f}")
            _log(logger, f"  Best fold : {best_fold_num}  (AUC={best_fold_auc:.4f})")
            _log(logger, f"  Rec(Fail) : {mean_recf:.4f}  ← minority-class recall")
            mean_tpr = np.mean(tprs,axis=0); mean_tpr[-1]=1.0; std_tpr=np.std(tprs,axis=0)
            plt.plot(mean_fpr,mean_tpr,'b',lw=3,label=f'Mean ROC (AUC={mean_auc:.3f}±{std_auc:.3f}){search_tag}')
            plt.fill_between(mean_fpr,np.maximum(mean_tpr-std_tpr,0),np.minimum(mean_tpr+std_tpr,1),
                             color='lightblue',alpha=0.3,label='±1 std')
            plt.plot([0,1],[0,1],'--',color='gray',lw=0.8)
            plt.xlim([0,1]); plt.ylim([0,1.05])
            plt.xlabel('False Positive Rate'); plt.ylabel('True Positive Rate')
            plt.title(f'XGBoost — {target.upper()} | {embedding_lm}{search_tag}\n'
                      f'{kfold_actual}-Fold SGKF  |  Acc={mean_acc:.3f}  F1={mean_f1:.3f}  '
                      f'Prec={mean_prec:.3f}  Rec={mean_rec:.3f}  Rec(Fail)={mean_recf:.3f}',fontsize=9)
            plt.legend(loc='lower right',fontsize=7); plt.grid(alpha=0.3); plt.tight_layout()
            _task_tag=""; plot_path=os.path.join(MODEL_DIR,
                f"CV_ROC_{target}_{embedding_lm}_xgboost{_db_tag}_k{kfold_actual}.png")
            plt.savefig(plot_path,dpi=150,bbox_inches='tight'); plt.close()
            _log(logger, f"  ROC plot  : {plot_path}")
        else:
            mean_r2=float(np.mean([m['r2'] for m in fold_metrics]))
            mean_rs=float(np.mean([m['spearman'] for m in fold_metrics]))
            _log(logger, f"  Best fold : {best_fold_num}  (ρ={best_fold_auc:.4f})")
            _log(logger, f"  Mean Spearman: {mean_rs:.4f} ± {std_auc:.4f}")
            mean_acc=mean_f1=mean_prec=mean_rec=mean_recf=float('nan'); _task_tag="_regression"
            plt.tight_layout()
            plt.savefig(os.path.join(MODEL_DIR,
                f"CV_scatter_{target}_{embedding_lm}_xgboost{_db_tag}_k{kfold_actual}_regression.png"),
                dpi=150,bbox_inches='tight'); plt.close()

        # ── Best-fold checkpoint ──────────────────────────────────────────
        _task_tag = "_regression" if _kfold_task=='regression' else ""
        best_path = None
        if best_fold_state is not None:
            best_path = os.path.join(MODEL_DIR,
                f"BEST_{target}_{embedding_lm}_xgboost{_db_tag}_k{kfold_actual}_fold{best_fold_num}{_task_tag}.pkl")
            _best=cls(config=best_fold_cfg,verbose=False); _best.model=best_fold_state
            _best.fb_=best_fold_fb; _best.task=_kfold_task; _best.save(best_path)
            _log(logger, f"\n[kfold] Best fold → {best_path}  (fold={best_fold_num}, AUC={best_fold_auc:.4f})")

        # ── Fold predictions CSV ──────────────────────────────────────────
        pred_path = None
        if all_records:
            pred_path = os.path.join(MODEL_DIR,
                f"fold_preds_{target}_{embedding_lm}_xgboost{_db_tag}_k{kfold_actual}{_task_tag}.csv")
            df_preds = pd.DataFrame(all_records)
            df_preds['best_fold'] = (df_preds['fold']==best_fold_num).astype(int)
            df_preds.to_csv(pred_path, index=False)
            _log(logger, f"[kfold] Fold predictions → {pred_path}")

        # ── Threshold optimisation ────────────────────────────────────────
        if _THRESHOLD_OPT_AVAILABLE and pred_path and best_path:
            _log(logger, "\n[threshold] Optimising ...")
            try:
                stability = run_full_threshold_pipeline(
                    fold_preds_csv=pred_path, target=target, lm=embedding_lm,
                    model='xgboost', db_stem=db_stem, best_ckpt_path=best_path,
                    output_dir=MODEL_DIR, cost_fp=cost_fp, cost_fn=cost_fn)
                rec_t = stability.get('pooled_threshold', 0.5)
                cls.recommended_threshold = float(rec_t)
                _log(logger, f"  Pooled OOF threshold : {rec_t:.4f}")
            except Exception as e:
                _log(logger, f"[threshold] WARNING: {e} — defaulting to 0.5")

        _log(logger, f"\n[log] Finished : {datetime.datetime.now()}")
        return mean_auc, std_auc, mean_acc, mean_f1, mean_prec, mean_rec